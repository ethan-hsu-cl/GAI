import json, logging, sys, re
from datetime import datetime
from pathlib import Path
from dataclasses import dataclass
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Cm, Inches
from pptx.enum.text import PP_ALIGN
from concurrent.futures import ThreadPoolExecutor

try: import cv2
except ImportError: cv2 = None

logging.basicConfig(level=logging.INFO, format='%(message)s')
logger = logging.getLogger(__name__)

@dataclass
class MediaPair:
    image_file: str
    image_path: Path
    video_path: Path = None
    ref_video_path: Path = None
    metadata: dict = None
    ref_metadata: dict = None
    failed: bool = False
    ref_failed: bool = False

def get_filename(folder, ref_folder='', model=''):
    def parse(name):
        m = re.match(r'(\d{4})\s+(.+)', name)
        return m.groups() if m else (datetime.now().strftime("%m%d"), name)
    d, s = parse(folder)
    parts = [f"[{d}]"]
    if model and model.lower() not in s.lower(): parts.append(model)
    if ref_folder:
        _, ref_s = parse(ref_folder)
        parts.append(f"{s} vs {ref_s}")
    else:
        parts.append(s)
    return ' '.join(parts)

class OptimizedVideoReportGenerator:
    def __init__(self, config_file="batch_config.json"):
        self.cfg = json.load(open(config_file, 'r', encoding='utf-8'))
        self._ar_cache = {}

    def get_ar(self, path, is_video=False):
        f = path.name.lower()
        if '9_16' in f or 'portrait' in f: return 9/16
        if '1_1' in f or 'square' in f: return 1
        if '16_9' in f or 'landscape' in f: return 16/9
        cache_key = str(path)
        if cache_key in self._ar_cache: return self._ar_cache[cache_key]
        try:
            if is_video and cv2:
                cap = cv2.VideoCapture(str(path))
                if cap.isOpened():
                    w, h = cap.get(3), cap.get(4)
                    cap.release()
                    if w > 0 and h > 0: 
                        self._ar_cache[cache_key] = w/h
                        return w/h
            else:
                with Image.open(path) as img: 
                    ar = img.width / img.height
                    self._ar_cache[cache_key] = ar
                    return ar
        except: pass
        return 16/9

    def add_media(self, slide, ph, media=None, is_video=False, poster=None, error=None):
        l, t, w, h = ph.left, ph.top, ph.width, ph.height
        ph._element.getparent().remove(ph._element)
        if media:
            try:
                ar = self.get_ar(Path(media), is_video)
                sw, sh = (w, w/ar) if ar > w/h else (h*ar, h)
                final_l, final_t = l + (w-sw)/2, t + (h-sh)/2
            except:
                sw, sh = w*0.8, h*0.8
                final_l, final_t = l + w*0.1, t + h*0.1
            if is_video:
                if poster and Path(poster).exists():
                    slide.shapes.add_movie(media, final_l, final_t, sw, sh, poster_frame_image=poster)
                else:
                    slide.shapes.add_movie(media, final_l, final_t, sw, sh)
            else:
                slide.shapes.add_picture(media, final_l, final_t, sw, sh)
        else:
            box = slide.shapes.add_textbox(l, t, w, h)
            box.text_frame.text = f"❌ GENERATION FAILED\n\n{error or 'No media found'}"
            for p in box.text_frame.paragraphs:
                p.font.size, p.alignment, p.font.color.rgb = Inches(16/72), PP_ALIGN.CENTER, RGBColor(255, 0, 0)
            box.fill.solid()
            box.fill.fore_color.rgb, box.line.color.rgb, box.line.width = RGBColor(255,240,240), RGBColor(255,0,0), Inches(0.02)

    def add_links(self, ppt, task):
        if not ppt.slides: return
        info_box = next((s for s in ppt.slides[0].shapes if hasattr(s, 'text_frame') and s.text_frame.text and any(x in s.text_frame.text.lower() for x in ['design', 'testbed', 'source'])), None)
        if not info_box: info_box = ppt.slides[0].shapes.add_textbox(Cm(5), Cm(13), Cm(20), Cm(4))
        info_box.text_frame.clear()
        links = [("Design: ", "Link", task.get('design_link', '')), ("Testbed: ", self.cfg.get('testbed', 'http://192.168.4.3:8000/video_effect/'), self.cfg.get('testbed', 'http://192.168.4.3:8000/video_effect/')), ("Source + Video: ", "Link", task.get('source_video_link', ''))]
        for i, (pre, text, url) in enumerate(links):
            para = info_box.text_frame.paragraphs[0] if i == 0 else info_box.text_frame.add_paragraph()
            if url:
                para.clear()
                r1, r2 = para.add_run(), para.add_run()
                r1.text, r1.font.size = pre, Inches(24/72)
                r2.text, r2.font.size = text if "Testbed:" not in pre else url, Inches(24/72)
                r2.hyperlink.address = url
                para.alignment = PP_ALIGN.CENTER
            else:
                para.text, para.font.size, para.alignment = f"{pre}{text}", Inches(24/72), PP_ALIGN.CENTER

    def find_matching_video(self, base, video_files):
        if base in video_files: return video_files[base]
        for vname, vpath in video_files.items():
            if vname.startswith(base + '_'): return vpath
        return None

    def process_batch(self, task):
        fp, ref_fp = Path(task['folder']), Path(task.get('reference_folder', '')) if task.get('reference_folder') else None
        use_cmp = task.get('use_comparison_template', False)
        folders = {'source': fp/'Source', 'video': fp/'Generated_Video', 'metadata': fp/'Metadata'}
        ref_folders = {'video': ref_fp/'Generated_Video', 'metadata': ref_fp/'Metadata'} if ref_fp and use_cmp else {}
        if not folders['source'].exists(): return fp, [], task
        
        files = {'images': {f.stem: f for f in folders['source'].iterdir() if f.suffix.lower() in {'.jpg', '.jpeg', '.png', '.webp'}},
                'videos': {f.stem: f for f in folders['video'].iterdir() if f.suffix.lower() in {'.mp4', '.mov'}} if folders['video'].exists() else {},
                'metadata': {f.stem.replace('_metadata', ''): f for f in folders['metadata'].iterdir() if f.suffix.lower() == '.json'} if folders['metadata'].exists() else {}}
        
        ref_files = {'videos': {f.stem: f for f in ref_folders['video'].iterdir() if f.suffix.lower() in {'.mp4', '.mov'}} if ref_folders.get('video', Path()).exists() else {},
                    'metadata': {f.stem.replace('_metadata', ''): f for f in ref_folders['metadata'].iterdir() if f.suffix.lower() == '.json'} if ref_folders.get('metadata', Path()).exists() else {}}
        
        def load_meta(base, meta_files):
            if base in meta_files:
                try: return json.load(open(meta_files[base]))
                except: return {}
            return {}
        
        with ThreadPoolExecutor(max_workers=8) as ex:
            bases = list(files['images'].keys())
            metadata_futures = [ex.submit(load_meta, base, files['metadata']) for base in bases]
            ref_metadata_futures = [ex.submit(load_meta, base, ref_files['metadata']) for base in bases] if use_cmp else [ex.submit(lambda: {}) for _ in bases]
            
            pairs = []
            for i, base in enumerate(bases):
                img = files['images'][base]
                md = metadata_futures[i].result()
                ref_md = ref_metadata_futures[i].result() if use_cmp else {}
                vid = self.find_matching_video(base, files['videos'])
                ref_vid = self.find_matching_video(base, ref_files['videos']) if use_cmp else None
                pairs.append(MediaPair(img.name, img, vid, ref_vid, md, ref_md, not vid or not md.get('success', False), use_cmp and (not ref_vid or not ref_md.get('success', False))))
        
        return fp, pairs, task

    def create_slide(self, ppt, pair, num, loaded, use_cmp=False):
        if loaded and len(ppt.slides) >= 4:
            slide = ppt.slides.add_slide(ppt.slides[3].slide_layout)
            for p in slide.placeholders:
                if p.placeholder_format.type == 1:
                    title = f"Generation {num}: {pair.image_file}"
                    if pair.failed: title += " ❌ FAILED"
                    if use_cmp and pair.ref_failed: title += " (REF ❌)"
                    p.text = title
                    if (pair.failed or (use_cmp and pair.ref_failed)) and p.text_frame.paragraphs:
                        p.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 0, 0)
                    break
            phs = sorted([p for p in slide.placeholders if p.placeholder_format.type in {6, 7, 8, 13, 18, 19}], key=lambda x: getattr(x, 'left', 0))
            if use_cmp and len(phs) >= 3:
                try: phs[0].insert_picture(str(pair.image_path))
                except: self.add_media(slide, phs[0], str(pair.image_path))
                self.add_media(slide, phs[1], str(pair.video_path) if pair.video_path and pair.video_path.exists() else None, True, str(pair.image_path), pair.metadata.get('error', 'Video generation failed') if pair.metadata else 'Generation failed')
                self.add_media(slide, phs[2], str(pair.ref_video_path) if pair.ref_video_path and pair.ref_video_path.exists() else None, True, str(pair.image_path), pair.ref_metadata.get('error', 'Reference video not found') if pair.ref_metadata else 'Reference video not found')
            elif len(phs) >= 2:
                try: phs[0].insert_picture(str(pair.image_path))
                except: self.add_media(slide, phs[0], str(pair.image_path))
                self.add_media(slide, phs[1], str(pair.video_path) if pair.video_path and pair.video_path.exists() else None, True, str(pair.image_path), pair.metadata.get('error', 'Video generation failed') if pair.metadata else 'Generation failed')
        else:
            slide = ppt.slides.add_slide(ppt.slide_layouts[6])
            title = f"Generation {num}: {pair.image_file}" + (" ❌ FAILED" if pair.failed else "") + (" (REF ❌)" if use_cmp and pair.ref_failed else "")
            tb = slide.shapes.add_textbox(Cm(2), Cm(1), Cm(20), Cm(2))
            tb.text_frame.text = title
            tb.text_frame.paragraphs[0].font.size = Inches(20/72)
            if pair.failed or (use_cmp and pair.ref_failed): tb.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 0, 0)
            
            pos = [(2.59, 3.26), (13, 3.26), (23.41, 3.26)] if use_cmp else [(2.59, 3.26), (18.78, 3.26)]
            w, h = (10, 10) if use_cmp else (12.5, 12.5)
            
            slide.shapes.add_picture(str(pair.image_path), Cm(pos[0][0]), Cm(pos[0][1]), Cm(w), Cm(h))
            if pair.video_path and pair.video_path.exists():
                try: slide.shapes.add_movie(str(pair.video_path), Cm(pos[1][0]), Cm(pos[1][1]), Cm(w), Cm(h), poster_frame_image=str(pair.image_path))
                except: slide.shapes.add_movie(str(pair.video_path), Cm(pos[1][0]), Cm(pos[1][1]), Cm(w), Cm(h))
            else: self._add_err(slide, pos[1], (w, h), pair.metadata.get('error', 'Video generation failed') if pair.metadata else 'Generation failed')
            
            if use_cmp:
                if pair.ref_video_path and pair.ref_video_path.exists():
                    try: slide.shapes.add_movie(str(pair.ref_video_path), Cm(pos[2][0]), Cm(pos[2][1]), Cm(w), Cm(h), poster_frame_image=str(pair.image_path))
                    except: slide.shapes.add_movie(str(pair.ref_video_path), Cm(pos[2][0]), Cm(pos[2][1]), Cm(w), Cm(h))
                else: self._add_err(slide, pos[2], (w, h), pair.ref_metadata.get('error', 'Reference video not found') if pair.ref_metadata else 'Reference video not found')
        
        meta_text = f"Task: {pair.metadata.get('task_id', 'N/A') if pair.metadata else 'N/A'}, Time: {pair.metadata.get('processing_time_seconds', 'N/A') if pair.metadata else 'N/A'}s"
        if use_cmp: meta_text += f"\nRef Task: {pair.ref_metadata.get('task_id', 'N/A') if pair.ref_metadata else 'N/A'}, Time: {pair.ref_metadata.get('processing_time_seconds', 'N/A') if pair.ref_metadata else 'N/A'}s"
        mb = slide.shapes.add_textbox(Cm(5), Cm(16), Cm(15), Cm(3))
        mb.text_frame.text = meta_text
        for p in mb.text_frame.paragraphs: p.font.size = Inches(10/72)

    def _add_err(self, slide, pos, size, msg):
        box = slide.shapes.add_textbox(Cm(pos[0]), Cm(pos[1]), Cm(size[0]), Cm(size[1]))
        box.text_frame.text = f"❌ GENERATION FAILED\n\n{msg}"
        for p in box.text_frame.paragraphs: p.font.size, p.alignment, p.font.color.rgb = Inches(16/72), PP_ALIGN.CENTER, RGBColor(255, 0, 0)
        box.fill.solid()
        box.fill.fore_color.rgb, box.line.color.rgb, box.line.width = RGBColor(255,240,240), RGBColor(255,0,0), Inches(0.02)

    def create_presentation(self, fp, pairs, task):
        if not pairs: return False
        use_cmp = task.get('use_comparison_template', False)
        ref_folder = task.get('reference_folder', '')
        tpath = self.cfg.get('comparison_template_path' if use_cmp else 'template_path', 'I2V Comparison Template.pptx' if use_cmp else 'I2V templates.pptx')
        try: ppt, loaded = (Presentation(tpath), True) if Path(tpath).exists() else (Presentation(), False)
        except: ppt, loaded = Presentation(), False
        ppt.slide_width, ppt.slide_height = Cm(33.87), Cm(19.05)
        
        fname = Path(fp).name
        if ppt.slides:
            if use_cmp and ref_folder:
                m1, m2 = re.match(r'(\d{4})\s+(.+)', fname), re.match(r'(\d{4})\s+(.+)', Path(ref_folder).name)
                d, s1 = m1.groups() if m1 else (datetime.now().strftime("%m%d"), fname)
                _, s2 = m2.groups() if m2 else ('', Path(ref_folder).name)
                title = f"[{d}] Kling 2.1\n{s1} vs {s2}"
            else:
                m = re.match(r'(\d{4})\s+(.+)', fname)
                d, s = m.groups() if m else (datetime.now().strftime("%m%d"), fname)
                title = f"[{d}] Kling 2.1\n{s}"
            for shape in ppt.slides[0].shapes:
                if hasattr(shape, 'text_frame') and (not shape.text_frame.text or "Results" in shape.text_frame.text):
                    shape.text_frame.clear()
                    p1 = shape.text_frame.paragraphs[0]
                    p1.text, p1.font.size, p1.alignment = title, Inches(48/72), PP_ALIGN.CENTER
                    break
        
        self.add_links(ppt, task)
        for i, pair in enumerate(pairs, 1): self.create_slide(ppt, pair, i, loaded, use_cmp)
        
        filename = get_filename(fname, Path(ref_folder).name if use_cmp and ref_folder else '', "Kling 2.1")
        opath = Path(self.cfg.get('output_directory', './')) / f"{filename}.pptx"
        ppt.save(str(opath))
        logger.info(f"✓ Saved: {opath}")
        return True

    def run(self):
        with ThreadPoolExecutor(max_workers=4) as ex:
            results = [f.result() for f in [ex.submit(self.process_batch, t) for t in self.cfg.get('tasks', [])]]
        successful = sum(1 for fp, pairs, task in results if self.create_presentation(fp, pairs, task))
        logger.info(f"✓ Generated {successful}/{len(results)} presentations")
        return successful > 0

if __name__ == "__main__":
    sys.exit(0 if OptimizedVideoReportGenerator().run() else 1)
