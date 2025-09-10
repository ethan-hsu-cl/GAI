import json, re, sys
from datetime import datetime
from pathlib import Path
from dataclasses import dataclass
from typing import List
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Cm, Inches
from pptx.enum.text import PP_ALIGN
from concurrent.futures import ThreadPoolExecutor
import logging

logging.basicConfig(level=logging.INFO, format='%(message)s')
logger = logging.getLogger(__name__)

@dataclass
class MediaPair:
    source_file: str
    source_path: Path
    gen_paths: List[Path] = None
    ref_paths: List[Path] = None
    metadata: dict = None
    failed: bool = False

def get_filename(folder, model=''):
    m = re.match(r'(\d{4})\s+(.+)', Path(folder).name)
    d,s = m.groups() if m else (datetime.now().strftime("%m%d"), Path(folder).name)
    parts = [f"[{d}]"]
    if model and model.lower() not in s.lower(): parts.append(model)
    parts.append(s)
    return ' '.join(parts)

def get_cmp_filename(f1, f2, model=''):
    m1, m2 = re.match(r'(\d{4})\s+(.+)', Path(f1).name), re.match(r'(\d{4})\s+(.+)', Path(f2).name)
    d,s1 = m1.groups() if m1 else (datetime.now().strftime("%m%d"), Path(f1).name)
    _,s2 = m2.groups() if m2 else ('', Path(f2).name)
    parts = [f"[{d}]"]
    if model and model.lower() not in s1.lower(): parts.append(model)
    parts.append(f"{s1} vs {s2}")
    return ' '.join(parts)

class NanoBananaReport:
    def __init__(self, cfg="batch_nano_banana_config.json"):
        self.cfg = json.load(open(cfg, 'r', encoding='utf-8'))
        self.pos = {'source': (2.59,3.26,12.5,12.5), 'generated': (18.78,3.26,12.5,12.5), 'reference': (35,3.26,12.5,12.5)}
        self.exts = {'.jpg','.jpeg','.png','.webp'}

    def add_links(self, ppt, task):
        if not ppt.slides: return
        slide = ppt.slides[0]
        info_box = next((s for s in slide.shapes if hasattr(s, 'text_frame') and s.text_frame.text and any(x in s.text_frame.text.lower() for x in ['design','testbed','source'])), None)
        if not info_box: info_box = slide.shapes.add_textbox(Cm(5), Cm(13), Cm(20), Cm(4))
        info_box.text_frame.clear()
        testbed_url = self.cfg.get('testbed', 'http://192.168.4.3:8000/video_effect/')
        lines = ["Design: Link", f"Testbed: {testbed_url}", "Source + Video: Link"]
        for i,line in enumerate(lines):
            p = info_box.text_frame.paragraphs[0] if i==0 else info_box.text_frame.add_paragraph()
            p.text,p.font.size,p.alignment = line, Inches(24/72), PP_ALIGN.CENTER
            if "Design:" in line and task.get('design_link'):
                self._add_hyperlink(p, "Design: ", "Link", task['design_link'])
            elif "Testbed:" in line:
                self._add_hyperlink(p, "Testbed: ", testbed_url, testbed_url)
            elif "Source + Video:" in line and task.get('source_video_link'):
                self._add_hyperlink(p, "Source + Video: ", "Link", task['source_video_link'])

    def _add_hyperlink(self, para, pre, link_text, url):
        para.clear()
        r1 = para.add_run(); r1.text,r1.font.size = pre, Inches(24/72)
        r2 = para.add_run(); r2.text,r2.font.size = link_text, Inches(24/72)
        r2.hyperlink.address = url
        para.alignment = PP_ALIGN.CENTER

    def match_batch(self):
        def proc(task):
            fp = Path(task['folder']); ref = task.get('reference_folder',''); use_cmp = task.get('use_comparison_template',False)
            folders = {'source': fp/'Source', 'output': fp/'Generated_Output', 'metadata': fp/'Metadata'}
            if ref and use_cmp: folders['reference'] = Path(ref)/'Generated_Output'
            if not folders['source'].exists(): return fp, [], task
            src = {f.stem:f for f in folders['source'].iterdir() if f.suffix.lower() in self.exts}
            out, ref_files = {}, {}
            if folders['output'].exists():
                for f in folders['output'].iterdir():
                    if f.suffix.lower() in self.exts and '_image_' in f.name:
                        out.setdefault(f.name.split('_image_')[0], []).append(f)
            if 'reference' in folders and folders['reference'].exists():
                for f in folders['reference'].iterdir():
                    if f.suffix.lower() in self.exts and '_image_' in f.name:
                        ref_files.setdefault(f.name.split('_image_')[0], []).append(f)
            pairs = []
            for b in sorted(src.keys()):
                md_file = folders['metadata'] / f"{b}_metadata.json"
                md = {}
                if md_file.exists():
                    try: md=json.load(open(md_file))
                    except: pass
                gen = sorted(out.get(b, []), key=lambda x: x.name)
                rf = sorted(ref_files.get(b, []), key=lambda x: x.name)
                pairs.append(MediaPair(src[b].name, src[b], gen, rf, md, not gen or not md.get('success',False)))
            return fp,pairs,task
        with ThreadPoolExecutor(max_workers=4) as ex:
            results = {str(fp):{'pairs':pairs,'task':task} for fp,pairs,task in [f.result() for f in [ex.submit(proc,t) for t in self.cfg.get('tasks',[])] ]}
        return results

    def calc_pos(self, path,x,y,w,h):
        try:
            with Image.open(path) as img: ar = img.width/img.height
            bw,bh = Cm(w),Cm(h)
            sw,sh = (bw,bw/ar) if ar>bw/bh else (bh*ar,bh)
            return Cm(x)+(bw-sw)/2,Cm(y)+(bh-sh)/2,sw,sh
        except: return Cm(x+0.5),Cm(y+0.5),Cm(w-1),Cm(h-1)

    def add_media(self, slide, ph, media=None, error=None):
        l,t,w,h=ph.left,ph.top,ph.width,ph.height
        ph._element.getparent().remove(ph._element)
        if media:
            try:
                with Image.open(media) as img: ar=img.width/img.height
                sw,sh=(w,w/ar) if ar>w/h else (h*ar,h)
                slide.shapes.add_picture(str(media), l+(w-sw)/2,t+(h-sh)/2, sw, sh)
            except: slide.shapes.add_picture(str(media), l+w*0.1,t+h*0.1, w*0.8,h*0.8)
        else:
            box=slide.shapes.add_textbox(l,t,w,h)
            box.text_frame.text=f"❌ GENERATION FAILED\n\n{error}"
            for p in box.text_frame.paragraphs:
                p.font.size,p.alignment,p.font.color.rgb = Inches(16/72),PP_ALIGN.CENTER,RGBColor(255,0,0)
            box.fill.solid(); box.fill.fore_color.rgb, box.line.color.rgb, box.line.width=RGBColor(255,240,240),RGBColor(255,0,0), Inches(0.02)

    def create_slide(self, ppt, pair, loaded, use_cmp):
        if loaded and len(ppt.slides)>=4:
            slide=ppt.slides.add_slide(ppt.slides[3].slide_layout)
            for p in slide.placeholders:
                if p.placeholder_format.type==1:
                    p.text="❌ GENERATION FAILED" if pair.failed else ""
                    if pair.failed and p.text_frame.paragraphs: p.text_frame.paragraphs[0].font.color.rgb=RGBColor(255,0,0)
                    break
            phs=sorted([p for p in slide.placeholders if p.placeholder_format.type in {6,7,8,13,18,19}], key=lambda x:getattr(x,'left',0))
            if use_cmp and len(phs)>=3:
                self.add_media(slide, phs[0], str(pair.source_path))
                self.add_media(slide, phs[1], str(pair.gen_paths[0]) if pair.gen_paths else None, "No images generated")
                self.add_media(slide, phs[2], str(pair.ref_paths[0]) if pair.ref_paths else None, "No reference images")
            elif len(phs)>=2:
                self.add_media(slide, phs[0], str(pair.source_path))
                self.add_media(slide, phs[1], str(pair.gen_paths[0]) if pair.gen_paths else None, "No images generated")
        else:
            slide=ppt.slides.add_slide(ppt.slide_layouts[6])
            if pair.failed:
                tb=slide.shapes.add_textbox(Cm(2),Cm(1),Cm(20),Cm(2))
                tb.text_frame.text,tb.text_frame.paragraphs[0].font.size,tb.text_frame.paragraphs[0].font.color.rgb= "❌ GENERATION FAILED",Inches(18/72),RGBColor(255,0,0)
            slide.shapes.add_picture(str(pair.source_path),*self.calc_pos(pair.source_path,*self.pos['source']))
            if pair.gen_paths:
                slide.shapes.add_picture(str(pair.gen_paths[0]),*self.calc_pos(pair.gen_paths[0],*self.pos['generated']))
            else: self._add_err(slide,self.pos['generated'],"No images generated")
            if use_cmp:
                if pair.ref_paths:
                    slide.shapes.add_picture(str(pair.ref_paths[0]),*self.calc_pos(pair.ref_paths[0],*self.pos['reference']))
                else: self._add_err(slide,self.pos['reference'],"No reference images")
        mb=slide.shapes.add_textbox(Cm(5),Cm(16),Cm(12),Cm(3))
        mb.text_frame.text=f"File Name: {pair.source_file}\nResponse ID: {pair.metadata.get('response_id','N/A') if pair.metadata else 'N/A'}\nTime: {pair.metadata.get('processing_time_seconds','N/A') if pair.metadata else 'N/A'}s"
        for p in mb.text_frame.paragraphs: p.font.size=Inches(10/72)

    def _add_err(self, slide, pos, msg):
        box=slide.shapes.add_textbox(Cm(pos[0]),Cm(pos[1]),Cm(pos[2]),Cm(pos[3]))
        box.text_frame.text=f"❌ FAILED\n\n{msg}"
        for p in box.text_frame.paragraphs:
            p.font.size,p.alignment,p.font.color.rgb=Inches(14/72),PP_ALIGN.CENTER,RGBColor(255,0,0)
        box.fill.solid()
        box.fill.fore_color.rgb,box.line.color.rgb=RGBColor(255,240,240),RGBColor(255,0,0)

    def create_presentation(self, folder, pairs, task):
        if not pairs: return False
        use_cmp, ref_folder=task.get('use_comparison_template',False),task.get('reference_folder','')
        tpath=self.cfg.get('comparison_template_path' if use_cmp else 'template_path') or ('I2V Comparison Template.pptx' if use_cmp else 'I2V templates.pptx')
        try: ppt, loaded = (Presentation(tpath), True) if Path(tpath).exists() else (Presentation(), False)
        except: ppt, loaded = Presentation(), False
        ppt.slide_width,ppt.slide_height=Cm(33.87),Cm(19.05)
        fname=Path(folder).name
        if ppt.slides and ppt.slides[0].shapes:
            if use_cmp and ref_folder:
                m1,m2=re.match(r'(\d{4})\s+(.+)',fname),re.match(r'(\d{4})\s+(.+)',Path(ref_folder).name)
                d,s1=m1.groups() if m1 else (datetime.now().strftime("%m%d"),fname)
                _,s2=m2.groups() if m2 else ('',Path(ref_folder).name)
                title=f"[{d}] Nano Banana\n{s1} vs {s2}"
            else:
                m=re.match(r'(\d{4})\s+(.+)',fname)
                d,s=m.groups() if m else (datetime.now().strftime("%m%d"),fname)
                title=f"[{d}] Nano Banana\n{s}"
            ppt.slides[0].shapes[0].text_frame.text=title
        self.add_links(ppt, task)
        for p in pairs: self.create_slide(ppt,p,loaded,use_cmp)
        fname=get_cmp_filename(folder,ref_folder,"Nano Banana") if use_cmp and ref_folder else get_filename(fname,"Nano Banana")
        out_dir=Path(self.cfg.get('output',{}).get('directory','./'))
        opath=out_dir/f"{fname}.pptx"
        ppt.save(str(opath))
        logger.info(f"✓ Saved: {opath}")
        return True

    def run(self):
        results=self.match_batch()
        successful=sum(1 for _,v in results.items() if self.create_presentation(_,v['pairs'],v['task']))
        logger.info(f"✓ Generated {successful}/{len(results)} presentations")
        return successful>0

if __name__=='__main__':
    sys.exit(0 if NanoBananaReport().run() else 1)
