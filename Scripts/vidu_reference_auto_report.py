import json, logging, sys, re
from datetime import datetime
from pathlib import Path
from dataclasses import dataclass
from collections import defaultdict
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Cm, Inches
from pptx.enum.text import PP_ALIGN

try: import cv2
except ImportError: cv2 = None

logging.basicConfig(level=logging.INFO, format='%(message)s')
logger = logging.getLogger(__name__)

@dataclass
class MediaPair:
    image_file: str
    image_path: Path
    effect_name: str
    category: str = "Reference"
    video_file: str = None
    video_path: Path = None
    metadata: dict = None
    failed: bool = False

def get_report_filename(folder_name: str, model_name: str = '') -> str:
    match = re.match(r'(\d{4})\s+(.+)', folder_name)
    date, style = match.groups() if match else (datetime.now().strftime("%m%d"), folder_name)
    parts = [f"[{date}]"]
    if model_name and model_name.lower() not in style.lower(): parts.append(model_name)
    parts.append(style)
    return ' '.join(parts)

def normalize_key(name: str) -> str:
    return re.sub(r'[^a-z0-9_]', '', name.lower().replace(' ', '_')).strip('_')

class ViduReferenceReportGenerator:
    def __init__(self, config_file: str = "batch_vidu_reference_config.json"):
        with open(config_file, 'r', encoding='utf-8') as f: 
            self.config = json.load(f)
        self.positions = {'img': (2.59, 3.26, 12.5, 12.5), 'vid': (18.78, 3.26, 12.5, 12.5)}
    
    def get_aspect_ratio(self, path: Path, is_video=False) -> float:
        name = path.name.lower()
        if "9_16" in name or "portrait" in name: return 9/16
        if "1_1" in name or "square" in name: return 1
        if "16_9" in name or "landscape" in name: return 16/9
        
        try:
            if is_video and cv2:
                cap = cv2.VideoCapture(str(path))
                if cap.isOpened():
                    w, h = cap.get(3), cap.get(4)
                    cap.release()
                    if w > 0 and h > 0: return w/h
            else:
                with Image.open(path) as img: return img.width / img.height
        except: pass
        return 16/9
    
    def fit_media_in_placeholder(self, slide, placeholder, media_path, is_video=False, poster_image=None):
        p_left, p_top, p_width, p_height = placeholder.left, placeholder.top, placeholder.width, placeholder.height
        ar = self.get_aspect_ratio(Path(media_path), is_video)
        
        if ar > p_width / p_height:
            scaled_w, scaled_h = p_width, p_width / ar
        else:
            scaled_h, scaled_w = p_height, p_height * ar
        
        final_left, final_top = p_left + (p_width - scaled_w) / 2, p_top + (p_height - scaled_h) / 2
        placeholder._element.getparent().remove(placeholder._element)
        
        try:
            if is_video:
                if poster_image: 
                    slide.shapes.add_movie(str(media_path), final_left, final_top, scaled_w, scaled_h, poster_frame_image=str(poster_image))
                else: 
                    slide.shapes.add_movie(str(media_path), final_left, final_top, scaled_w, scaled_h)
            else:
                slide.shapes.add_picture(str(media_path), final_left, final_top, scaled_w, scaled_h)
        except:
            slide.shapes.add_movie(str(media_path), final_left, final_top, scaled_w, scaled_h) if is_video else slide.shapes.add_picture(str(media_path), final_left, final_top, scaled_w, scaled_h)
    
    def add_error_to_placeholder(self, slide, placeholder, error_msg):
        p_left, p_top, p_width, p_height = placeholder.left, placeholder.top, placeholder.width, placeholder.height
        placeholder._element.getparent().remove(placeholder._element)
        
        error_box = slide.shapes.add_textbox(p_left, p_top, p_width, p_height)
        error_box.text_frame.text = f"❌ REFERENCE FAILED\n\n{error_msg}"
        for p in error_box.text_frame.paragraphs:
            p.font.size, p.alignment, p.font.color.rgb = Inches(16/72), PP_ALIGN.CENTER, RGBColor(255, 0, 0)
        error_box.fill.solid()
        error_box.fill.fore_color.rgb, error_box.line.color.rgb = RGBColor(255, 240, 240), RGBColor(255, 0, 0)
        error_box.line.width = Inches(0.02)
    
    def extract_key(self, filename: str, effect: str) -> str:
        stem = Path(filename).stem
        effect_clean = normalize_key(effect)
        
        if stem.lower().endswith(f'_{effect_clean}'):
            key = stem[:-len(f'_{effect_clean}')]
        else:
            for suffix in ['_reference', '_generated', '_output', '_result']:
                if stem.lower().endswith(suffix):
                    key = stem[:-len(suffix)]
                    break
            else:
                key = stem
        return normalize_key(key)
    
    def discover_folders(self):
        base = Path(self.config.get('base_folder', ''))
        if not base.exists(): raise FileNotFoundError(f"Base folder not found: {base}")
        return sorted(f.name for f in base.iterdir() 
                     if f.is_dir() and not f.name.startswith(('.', '_')) 
                     and (f / 'Source').exists() and (f / 'Reference').exists())
    
    def process_folders(self):
        base_folder = Path(self.config.get('base_folder', ''))
        
        try:
            effect_names = self.discover_folders()
            logger.info(f"🔍 Discovered {len(effect_names)} effect folders")
        except:
            effect_names = [t.get('effect', '') for t in self.config.get('tasks', [])]
            logger.info(f"📋 Using {len(effect_names)} configured tasks")
        
        pairs = []
        for effect in effect_names:
            if not effect: continue
            
            folders = {k: base_folder / effect / v for k, v in 
                      {'src': 'Source', 'vid': 'Generated_Video', 'meta': 'Metadata'}.items()}
            
            if not folders['src'].exists(): continue
            logger.info(f"Processing: {effect}")
            
            # Get images and videos
            images = {normalize_key(f.stem): f for f in folders['src'].iterdir() 
                     if f.suffix.lower() in {'.jpg', '.jpeg', '.png', '.webp'}}
            
            videos = {}
            if folders['vid'].exists():
                for f in folders['vid'].iterdir():
                    if f.suffix.lower() in {'.mp4', '.mov', '.avi'}:
                        videos[self.extract_key(f.name, effect)] = f
            
            logger.info(f"📊 Images: {len(images)}, Videos: {len(videos)}")
            
            # Create pairs
            for key, img in images.items():
                meta_file = folders['meta'] / f"{img.stem}_metadata.json"
                metadata = {}
                if meta_file.exists():
                    try:
                        with open(meta_file, 'r', encoding='utf-8') as f: 
                            metadata = json.load(f)
                    except: pass
                
                vid = videos.get(key)
                logger.info(f"  {'✓' if vid else '❌'} Matched: {key}")
                
                pairs.append(MediaPair(
                    img.name, img, effect, "Reference", 
                    vid.name if vid else None, vid, metadata,
                    not vid or not metadata.get('success', False)
                ))
        
        return pairs
    
    def create_slide(self, ppt, pair, idx):
        # Try template layout first
        template_loaded = False
        try:
            if len(ppt.slides) >= 4:
                slide = ppt.slides.add_slide(ppt.slides[3].slide_layout)
                template_loaded = True
            else:
                slide = ppt.slides.add_slide(ppt.slide_layouts[6])
        except:
            slide = ppt.slides.add_slide(ppt.slide_layouts[6])
        
        ref_count = pair.metadata.get('reference_count', 0) if pair.metadata else 0
        title = f"{pair.effect_name} #{idx}: {pair.image_file} (+{ref_count} refs)"
        if pair.failed: title += " ❌ FAILED"
        
        if template_loaded:
            placeholders = sorted([p for p in slide.placeholders], 
                                key=lambda x: x.left if hasattr(x, 'left') else 0)
            
            # Update title
            for p in slide.placeholders:
                if p.placeholder_format.type == 1:
                    p.text = title
                    if pair.failed and p.text_frame.paragraphs:
                        p.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 0, 0)
                    break
            
            # Handle content placeholders
            content_pls = [p for p in placeholders if p.placeholder_format.type in {6, 7, 8, 13, 18, 19}]
            if len(content_pls) >= 2:
                self.fit_media_in_placeholder(slide, content_pls[0], str(pair.image_path), False)
                if pair.video_path and pair.video_path.exists():
                    self.fit_media_in_placeholder(slide, content_pls[1], str(pair.video_path), 
                                                True, str(pair.image_path))
                else:
                    error_msg = pair.metadata.get('error_message', 'Reference generation failed') if pair.metadata else 'No processing'
                    self.add_error_to_placeholder(slide, content_pls[1], error_msg)
        
        else:
            # Manual positioning
            title_box = slide.shapes.add_textbox(Cm(1), Cm(0.5), Cm(30), Cm(2))
            title_box.text_frame.text = title
            title_box.text_frame.paragraphs[0].font.size = Inches(18/72)
            if pair.failed: title_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 0, 0)
            
            # Add image
            img_x, img_y, img_w, img_h = self.calc_position(pair.image_path, *self.positions['img'])
            slide.shapes.add_picture(str(pair.image_path), img_x, img_y, img_w, img_h)
            
            # Add video or error
            if pair.video_path and pair.video_path.exists():
                vid_x, vid_y, vid_w, vid_h = self.calc_position(pair.video_path, *self.positions['vid'], True)
                try: 
                    slide.shapes.add_movie(str(pair.video_path), vid_x, vid_y, vid_w, vid_h, 
                                         poster_frame_image=str(pair.image_path))
                except: 
                    slide.shapes.add_movie(str(pair.video_path), vid_x, vid_y, vid_w, vid_h)
            else:
                pos = self.positions['vid']
                error_box = slide.shapes.add_textbox(Cm(pos[0]), Cm(pos[1]), Cm(pos[2]), Cm(pos[3]))
                error_msg = pair.metadata.get('error_message', 'Reference generation failed') if pair.metadata else 'No processing'
                error_box.text_frame.text = f"❌ REFERENCE FAILED\n\n{error_msg}"
                for p in error_box.text_frame.paragraphs:
                    p.font.size, p.alignment, p.font.color.rgb = Inches(14/72), PP_ALIGN.CENTER, RGBColor(255, 0, 0)
                error_box.fill.solid()
                error_box.fill.fore_color.rgb, error_box.line.color.rgb = RGBColor(255, 240, 240), RGBColor(255, 0, 0)
        
        # Add metadata box
        proc_time = pair.metadata.get('processing_time_seconds', 'N/A') if pair.metadata else 'N/A'
        aspect_ratio = pair.metadata.get('detected_aspect_ratio', 'N/A') if pair.metadata else 'N/A'
        status = 'FAILED' if pair.failed else 'SUCCESS'
        
        meta_text = f"{pair.image_file}\n{pair.effect_name} | +{ref_count} refs\n{aspect_ratio} | {proc_time}s\n{status}"
        meta_box = slide.shapes.add_textbox(Cm(2), Cm(16.5), Cm(10), Cm(2.5))
        meta_box.text_frame.text = meta_text
        meta_box.text_frame.word_wrap = True
        for p in meta_box.text_frame.paragraphs: 
            p.font.size, p.alignment = Inches(9/72), PP_ALIGN.LEFT
    
    def calc_position(self, path, x, y, w, h, is_video=False):
        try:
            ar = self.get_aspect_ratio(path, is_video)
            box_w, box_h = Cm(w), Cm(h)
            if ar > box_w / box_h:
                scaled_w, scaled_h = box_w, box_w / ar
            else:
                scaled_h, scaled_w = box_h, box_h * ar
            return Cm(x) + (box_w - scaled_w) / 2, Cm(y) + (box_h - scaled_h) / 2, scaled_w, scaled_h
        except:
            return Cm(x + 0.5), Cm(y + 0.5), Cm(w - 1), Cm(h - 1)
    
    def run(self):
        try:
            pairs = self.process_folders()
            if not pairs: 
                logger.warning("No media pairs found")
                return False
            
            # Create presentation
            try: 
                ppt = Presentation(self.config.get('template_path', 'I2V templates.pptx'))
            except: 
                ppt = Presentation()
            ppt.slide_width, ppt.slide_height = Cm(33.87), Cm(19.05)
            
            # Title slide
            base_name = Path(self.config.get('base_folder', '')).name
            match = re.match(r'(\d{4})\s+(.+)', base_name)
            date, project = match.groups() if match else (datetime.now().strftime("%m%d"), base_name)
            
            if ppt.slides and ppt.slides[0].shapes:
                ppt.slides[0].shapes[0].text_frame.text = f"[{date}] Vidu Reference\n{project}"
            
            # Group pairs by style and create slides
            grouped = defaultdict(list)
            for p in pairs:
                grouped[p.effect_name].append(p)
            
            total_styles = len(grouped)
            
            for style_name, style_pairs in sorted(grouped.items()):
                
                # **UPDATED: Insert section slide using Title and Content layout**
                try:
                    # Use Title and Content layout (typically layout 1)
                    section_slide = ppt.slides.add_slide(ppt.slide_layouts[1])
                except:
                    # Fallback to blank layout if Title and Content not available
                    section_slide = ppt.slides.add_slide(ppt.slide_layouts[6])
                
                # Find and set the title placeholder
                for placeholder in section_slide.placeholders:
                    if placeholder.placeholder_format.type == 1:  # Title placeholder
                        placeholder.text = style_name  # Just the style name, no "Style:" prefix
                        break
                
                logger.info(f"📄 Section: {style_name}")
                
                # Add content slides for this style
                for idx, pair in enumerate(style_pairs, 1):
                    self.create_slide(ppt, pair, idx)
                    logger.info(f"📄 Content: {style_name} #{idx}")
            
            # Save presentation
            filename = get_report_filename(base_name, "Vidu Reference")
            output_path = Path(self.config.get('output_directory', './')) / f"{filename}.pptx"
            ppt.save(str(output_path))
            
            logger.info(f"✓ Saved: {output_path}")
            logger.info(f"📊 Summary: {total_styles} styles, {len(pairs)} total references")
            return True
            
        except Exception as e:
            logger.error(f"Report generation failed: {e}")
            return False

if __name__ == "__main__":
    sys.exit(0 if ViduReferenceReportGenerator().run() else 1)
