"""
Veo Report Generator - Text-to-video report generation.

Differences from other reports:
- No source images (text-to-video only)
- Prompt text displayed in place of source image
- Only generated video and metadata shown
"""
import json
import logging
import sys
from datetime import datetime
from pathlib import Path
from dataclasses import dataclass, field
from typing import Dict, List, Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Cm, Pt
from pptx.enum.text import PP_ALIGN

# Add parent directory to path for imports
sys.path.insert(0, str(Path(__file__).parent.parent))

logging.basicConfig(level=logging.INFO, format='%(message)s')
logger = logging.getLogger(__name__)


@dataclass
class VeoMediaPair:
    """
    Media pair for Veo text-to-video generation.
    
    Unlike other APIs, Veo has no source image - only text prompt.
    """
    prompt: str
    prompt_truncated: str  # Shortened version for display
    style_name: str
    generation_number: int
    
    generated_video_path: Optional[Path] = None
    metadata: Dict = field(default_factory=dict)
    failed: bool = False
    
    @property
    def display_name(self) -> str:
        """Get display name for this generation."""
        return f"{self.style_name}-{self.generation_number}"


class VeoReportGenerator:
    """
    Veo-specific report generator.
    
    Key differences from standard report generator:
    - No source image handling
    - Prompt text displayed as primary content
    - Single-media layout (only generated video)
    """
    
    IMAGE_EXTS = {'.jpg', '.jpeg', '.png', '.webp'}
    VIDEO_EXTS = {'.mp4', '.mov', '.avi'}
    
    # Layout positioning (matching unified_report_generator.py LAYOUT_2_MEDIA)
    LAYOUT_2_MEDIA = {
        'positions': [(0.42, 2.15, 16, 16), (17.44, 2.15, 16, 16)],
        'metadata_position': (35, 0, 7.29, 3.06),
    }
    
    def __init__(self, config_file: str = None):
        """
        Initialize Veo report generator.
        
        Args:
            config_file: Path to Veo configuration file
        """
        self.config_file = config_file or "config/batch_veo_config.json"
        self.config = {}
        
        self.load_config()
    
    def load_config(self):
        """Load Veo configuration."""
        try:
            with open(self.config_file, 'r', encoding='utf-8') as f:
                self.config = json.load(f)
            logger.info(f"✓ Config loaded: {self.config_file}")
        except Exception as e:
            logger.error(f"❌ Config error: {e}")
            sys.exit(1)
    
    def get_aspect_ratio(self, media_path: Path) -> float:
        """Get aspect ratio of video file."""
        try:
            import cv2
            cap = cv2.VideoCapture(str(media_path))
            if cap.isOpened():
                width = cap.get(cv2.CAP_PROP_FRAME_WIDTH)
                height = cap.get(cv2.CAP_PROP_FRAME_HEIGHT)
                cap.release()
                if height > 0:
                    return width / height
        except Exception as e:
            logger.warning(f"Failed to get aspect ratio for {media_path.name}: {e}")
        return 1.0  # Default to square
    
    def extract_first_frame(self, video_path: Path) -> Optional[str]:
        """Extract first frame from video for poster image."""
        try:
            import cv2
            import tempfile
            
            cap = cv2.VideoCapture(str(video_path))
            if not cap.isOpened():
                return None
            
            ret, frame = cap.read()
            cap.release()
            
            if ret:
                tmp = tempfile.NamedTemporaryFile(suffix='.jpg', delete=False)
                cv2.imwrite(tmp.name, frame)
                tmp.close()
                return tmp.name
        except Exception as e:
            logger.warning(f"Failed to extract frame from {video_path.name}: {e}")
        return None
    
    def process_batch(self) -> List[VeoMediaPair]:
        """
        Process all Veo tasks and create media pairs.
        
        Returns:
            List of VeoMediaPair objects
        """
        pairs = []
        tasks = self.config.get('tasks', [])
        
        global_count = self.config.get('generation_count', 1)
        
        for task in tasks:
            output_folder = Path(task.get('output_folder', ''))
            if not output_folder.exists():
                logger.warning(f"Output folder not found: {output_folder}")
                continue
            
            metadata_folder = output_folder.parent / "Metadata"
            if not metadata_folder.exists():
                logger.warning(f"Metadata folder not found: {metadata_folder}")
                continue
            
            # Get task details
            prompt = task.get('prompt', '')
            style_name = task.get('style_name', 'VeoTask')
            task_count = task.get('generation_count')
            generation_count = task_count if task_count is not None else global_count
            
            if generation_count < 1:
                generation_count = 1
            
            # Truncate prompt for display (first 100 chars)
            prompt_truncated = prompt[:100] + '...' if len(prompt) > 100 else prompt
            
            # Create safe filename from style name
            safe_style = "".join(c if c.isalnum() or c in (' ', '-', '_') else '_' for c in style_name)
            safe_style = safe_style.strip().replace(' ', '_')
            
            # Find all generations for this task
            for gen_num in range(1, generation_count + 1):
                base_name = f"{safe_style}-{gen_num}"
                
                # Find generated video
                video_path = output_folder / f"{base_name}_generated.mp4"
                if not video_path.exists():
                    # Try alternative extensions
                    for ext in ['.mov', '.avi']:
                        alt_path = output_folder / f"{base_name}_generated{ext}"
                        if alt_path.exists():
                            video_path = alt_path
                            break
                
                # Load metadata
                metadata_path = metadata_folder / f"{base_name}_metadata.json"
                metadata = {}
                if metadata_path.exists():
                    try:
                        with open(metadata_path, 'r', encoding='utf-8') as f:
                            metadata = json.load(f)
                    except Exception as e:
                        logger.warning(f"Failed to load metadata {metadata_path.name}: {e}")
                
                # Create pair
                pair = VeoMediaPair(
                    prompt=prompt,
                    prompt_truncated=prompt_truncated,
                    style_name=style_name,
                    generation_number=gen_num,
                    generated_video_path=video_path if video_path.exists() else None,
                    metadata=metadata,
                    failed=not video_path.exists() or not metadata.get('success', False)
                )
                pairs.append(pair)
                
                if pair.failed:
                    logger.warning(f"Failed generation: {base_name}")
                else:
                    logger.info(f"Valid generation: {base_name}")
        
        logger.info(f"Created {len(pairs)} Veo media pairs")
        return pairs
    
    def create_presentation(self, pairs: List[VeoMediaPair]) -> bool:
        """
        Create PowerPoint presentation for Veo results.
        
        Args:
            pairs: List of VeoMediaPair objects
            
        Returns:
            True if successful, False otherwise
        """
        if not pairs:
            logger.warning("No media pairs to process")
            return False
        
        # Load template
        template_path = self.config.get('template_path', 'templates/I2V templates.pptx')
        
        try:
            ppt = Presentation(template_path) if Path(template_path).exists() else Presentation()
            template_loaded = Path(template_path).exists()
            logger.info(f"✓ Template loaded: {template_path}")
        except Exception as e:
            logger.warning(f"⚠️ Template load failed: {e}, using blank presentation")
            ppt = Presentation()
            template_loaded = False
        
        # Set slide dimensions
        ppt.slide_width, ppt.slide_height = Cm(33.87), Cm(19.05)
        
        # Create title slide
        self.create_title_slide(ppt)
        
        # Create content slides
        for i, pair in enumerate(pairs, 1):
            self.create_content_slide(ppt, pair, i, template_loaded)
        
        # Save presentation
        return self.save_presentation(ppt, pairs)
    
    def create_title_slide(self, ppt: Presentation):
        """Create title slide for Veo report."""
        if not ppt.slides:
            return
        
        # Generate title
        date_str = datetime.now().strftime("%m%d")
        title = f"[{date_str}] Google Veo Text-to-Video"
        
        # Update title slide
        if ppt.slides and ppt.slides[0].shapes:
            ppt.slides[0].shapes[0].text_frame.text = title
        
        # Add testbed link
        slide = ppt.slides[0]
        info_box = next((s for s in slide.shapes if hasattr(s, 'text_frame') and s.text_frame.text and
                        any(k in s.text_frame.text.lower() for k in ['testbed', 'source'])), None)
        
        if not info_box:
            info_box = slide.shapes.add_textbox(Cm(5), Cm(13), Cm(20), Cm(4))
        
        info_box.text_frame.clear()
        
        # Add testbed link
        testbed_url = "http://192.168.4.8:8000/google_veo/"
        para = info_box.text_frame.paragraphs[0]
        para.clear()
        r1, r2 = para.add_run(), para.add_run()
        r1.text, r1.font.size = "Testbed: ", Pt(24)
        r2.text, r2.font.size = testbed_url, Pt(24)
        r2.hyperlink.address = testbed_url
        para.alignment = PP_ALIGN.CENTER
    
    def create_content_slide(self, ppt: Presentation, pair: VeoMediaPair, 
                            index: int, template_loaded: bool):
        """
        Create content slide for Veo generation.
        
        Layout:
        - Left side: Prompt text box
        - Right side: Generated video
        - Bottom: Metadata
        """
        # Create slide
        if template_loaded and len(ppt.slides) >= 4:
            slide = ppt.slides.add_slide(ppt.slides[3].slide_layout)
            self._handle_template_slide(slide, pair, index)
        else:
            slide = ppt.slides.add_slide(ppt.slide_layouts[6])
            self._handle_manual_slide(slide, pair, index)
    
    def _handle_template_slide(self, slide, pair: VeoMediaPair, index: int):
        """Handle slide with template placeholders."""
        # Update title if present
        title_ph = next((p for p in slide.placeholders if p.placeholder_format.type == 1), None)
        if title_ph:
            title_ph.text = f"Generation {index}: {pair.display_name}"
            if pair.failed and title_ph.text_frame.paragraphs:
                title_ph.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 0, 0)
        
        # Get media placeholders (sorted left to right)
        phs = sorted([p for p in slide.placeholders 
                     if p.placeholder_format.type in {6, 7, 8, 13, 18, 19}],
                    key=lambda x: getattr(x, 'left', 0))
        
        if len(phs) >= 2:
            # Left placeholder: Prompt text
            self._add_prompt_box(slide, phs[0], pair)
            
            # Right placeholder: Generated video
            self._add_video(slide, phs[1], pair)
        
        # Add metadata
        self._add_metadata(slide, pair)
    
    def _handle_manual_slide(self, slide, pair: VeoMediaPair, index: int):
        """Handle slide without template."""
        # Add title
        title_text = f"Generation {index}: {pair.display_name}"
        if pair.failed:
            title_text = f"❌ {title_text} - GENERATION FAILED"
        
        tb = slide.shapes.add_textbox(Cm(2), Cm(1), Cm(20), Cm(2))
        tb.text_frame.text = title_text
        tb.text_frame.paragraphs[0].font.size = Pt(20)
        if pair.failed:
            tb.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 0, 0)
        
        # Use LAYOUT_2_MEDIA positions
        positions = self.LAYOUT_2_MEDIA['positions']
        
        # Left side: Prompt (position like source image)
        self._add_prompt_box(slide, positions[0], pair)
        
        # Right side: Generated video
        self._add_video(slide, positions[1], pair)
        
        # Add metadata
        self._add_metadata(slide, pair)
    
    def _add_prompt_box(self, slide, placeholder_or_pos, pair: VeoMediaPair):
        """Add prompt text box in place of source image."""
        # Handle both placeholder objects and manual positions
        if hasattr(placeholder_or_pos, 'left'):
            l, t, w, h = (placeholder_or_pos.left, placeholder_or_pos.top,
                         placeholder_or_pos.width, placeholder_or_pos.height)
            try:
                placeholder_or_pos._element.getparent().remove(placeholder_or_pos._element)
            except:
                pass
        else:
            l, t, w, h = [Cm(x) for x in placeholder_or_pos]
        
        # Create text box for prompt
        box = slide.shapes.add_textbox(l, t, w, h)
        tf = box.text_frame
        tf.word_wrap = True
        
        # Add header
        header = tf.paragraphs[0]
        header.text = "TEXT PROMPT:"
        header.font.bold = True
        header.font.size = Pt(18)
        header.alignment = PP_ALIGN.CENTER
        
        # Add prompt text
        prompt_para = tf.add_paragraph()
        prompt_para.text = pair.prompt
        prompt_para.font.size = Pt(12)
        prompt_para.space_before = Pt(12)
        
        # Style the box
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(245, 245, 255)  # Light blue background
        box.line.color.rgb = RGBColor(100, 100, 200)
        box.line.width = Pt(2)
    
    def _add_video(self, slide, placeholder_or_pos, pair: VeoMediaPair):
        """Add generated video or error placeholder with proper aspect ratio."""
        # Handle both placeholder objects and manual positions
        if hasattr(placeholder_or_pos, 'left'):
            l, t, w, h = (placeholder_or_pos.left, placeholder_or_pos.top,
                         placeholder_or_pos.width, placeholder_or_pos.height)
            try:
                placeholder_or_pos._element.getparent().remove(placeholder_or_pos._element)
            except:
                pass
        else:
            l, t, w, h = [Cm(x) for x in placeholder_or_pos]
        
        if pair.generated_video_path and pair.generated_video_path.exists():
            # Add video with aspect ratio calculation
            try:
                # Calculate aspect ratio and positioning
                ar = self.get_aspect_ratio(pair.generated_video_path)
                sw, sh = (w, w/ar) if ar > w/h else (h*ar, h)
                fl, ft = l + (w - sw)/2, t + (h - sh)/2
                
                # Extract first frame for poster
                first_frame_path = self.extract_first_frame(pair.generated_video_path)
                if first_frame_path:
                    slide.shapes.add_movie(str(pair.generated_video_path), fl, ft, sw, sh,
                                         poster_frame_image=first_frame_path)
                else:
                    slide.shapes.add_movie(str(pair.generated_video_path), fl, ft, sw, sh)
            except Exception as e:
                self._add_error_box(slide, l, t, w, h, f"Failed to load video: {e}")
        else:
            # Add error box
            error_msg = pair.metadata.get('error', 'Video not found')
            status_msg = pair.metadata.get('status_message', '')
            full_msg = f"{error_msg}\n\n{status_msg}" if status_msg else error_msg
            self._add_error_box(slide, l, t, w, h, full_msg)
    
    def _add_error_box(self, slide, left, top, width, height, message: str):
        """Add error box for failed generation."""
        box = slide.shapes.add_textbox(left, top, width, height)
        box.text_frame.text = f"❌ GENERATION FAILED\n\n{message}"
        box.text_frame.word_wrap = True
        
        for para in box.text_frame.paragraphs:
            para.font.size = Pt(12)
            para.alignment = PP_ALIGN.CENTER
            para.font.color.rgb = RGBColor(255, 0, 0)
        
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(255, 240, 240)
        box.line.color.rgb = RGBColor(255, 0, 0)
        box.line.width = Pt(0.5)
    
    def _add_metadata(self, slide, pair: VeoMediaPair):
        """Add metadata text box."""
        md = pair.metadata
        
        meta_lines = [
            f"Style: {pair.style_name}",
            f"Generation: {pair.generation_number}",
            f"Status: {'✓ Success' if md.get('success', False) else '❌ Failed'}",
        ]
        
        # if md.get('model_id'):
        #     meta_lines.append(f"Model: {md['model_id']}")
        if md.get('duration_seconds'):
            meta_lines.append(f"Duration: {md['duration_seconds']}s")
        # if md.get('aspect_ratio'):
        #     meta_lines.append(f"Aspect Ratio: {md['aspect_ratio']}")
        # if md.get('resolution'):
        #     meta_lines.append(f"Resolution: {md['resolution']}")
        if md.get('processing_time_seconds'):
            meta_lines.append(f"Processing Time: {md['processing_time_seconds']}s")
        if md.get('attempts'):
            meta_lines.append(f"Attempts: {md['attempts']}")
        
        # Use LAYOUT_2_MEDIA metadata position
        metadata_pos = self.LAYOUT_2_MEDIA['metadata_position']
        box = slide.shapes.add_textbox(Cm(metadata_pos[0]), Cm(metadata_pos[1]),
                                       Cm(metadata_pos[2]), Cm(metadata_pos[3]))
        box.text_frame.text = "\n".join(meta_lines)
        box.text_frame.word_wrap = True
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(255, 255, 255)
        
        for para in box.text_frame.paragraphs:
            para.font.size = Pt(10)
    
    def save_presentation(self, ppt: Presentation, pairs: List[VeoMediaPair]) -> bool:
        """Save the presentation."""
        try:
            # Generate filename
            date_str = datetime.now().strftime("%m%d")
            
            # Collect unique style names
            style_names = []
            seen = set()
            for pair in pairs:
                if pair.style_name not in seen:
                    style_names.append(pair.style_name)
                    seen.add(pair.style_name)
            
            style_str = ', '.join(style_names) if style_names else 'Test'
            filename = f"[{date_str}] Google Veo {style_str}"
            
            # Get output directory
            output_dir = Path(self.config.get('output_directory', './'))
            output_dir.mkdir(parents=True, exist_ok=True)
            
            output_path = output_dir / f"{filename}.pptx"
            ppt.save(str(output_path))
            
            logger.info(f"✓ Saved: {output_path}")
            return True
        except Exception as e:
            logger.error(f"❌ Save failed: {e}")
            return False
    
    def run(self) -> bool:
        """Main execution."""
        logger.info("🎬 Starting Veo Report Generator")
        
        try:
            # Process batch
            pairs = self.process_batch()
            
            if not pairs:
                logger.warning("No media pairs found")
                return False
            
            # Create presentation
            return self.create_presentation(pairs)
            
        except Exception as e:
            logger.error(f"❌ Report generation failed: {e}")
            return False


def main():
    """CLI entry point."""
    import argparse
    
    parser = argparse.ArgumentParser(description='Generate PowerPoint report for Veo text-to-video results')
    parser.add_argument('--config', '-c', help='Config file path (optional)')
    
    args = parser.parse_args()
    
    generator = VeoReportGenerator(args.config)
    sys.exit(0 if generator.run() else 1)


if __name__ == "__main__":
    main()
