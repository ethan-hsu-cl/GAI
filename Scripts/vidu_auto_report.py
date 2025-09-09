import json, logging, sys, re
from datetime import datetime
from pathlib import Path
from dataclasses import dataclass
from typing import Dict
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Cm, Inches
from pptx.enum.text import PP_ALIGN

try: import cv2
except ImportError: cv2 = None

logging.basicConfig(level=logging.INFO, format='%(message)s')
logger = logging.getLogger(__name__)

@dataclass
class MediaPair:
    image_file: str
    image_path: Path
    effect_name: str
    category: str
    video_file: str = None
    video_path: Path = None
    metadata: Dict = None
    failed: bool = False

def get_report_filename(folder_name: str, model_name: str = '') -> str:
    match = re.match(r'(\d{4})\s+(.+)', folder_name)
    date, style = match.groups() if match else (datetime.now().strftime("%m%d"), folder_name)
    parts = [f"[{date}]"]
    if model_name and model_name.lower() not in style.lower(): parts.append(model_name)
    parts.append(style)
    return ' '.join(parts)

def normalize_key(name: str) -> str:
    """FIXED: Normalize key and strip trailing underscores"""
    key = name.lower().replace(' ', '_')
    key = re.sub(r'[^a-z0-9_]', '', key)
    key = re.sub(r'_effect$', '', key)
    return key.strip('_')  # **CRITICAL FIX for trailing underscores**

class ViduReportGenerator:
    def __init__(self, config_file: str = "batch_vidu_config.json"):
        with open(config_file, 'r') as f: self.config = json.load(f)
        self.positions = {'img': (2.59, 3.26, 12.5, 12.5), 'vid': (18.78, 3.26, 12.5, 12.5)}
    
    def get_aspect_ratio(self, path: Path, is_video=False) -> float:
        filename = path.name.lower()
        if "9_16" in filename or "portrait" in filename: return 9/16
        if "1_1" in filename or "square" in filename: return 1
        if "16_9" in filename or "landscape" in filename: return 16/9
        
        try:
            if is_video and cv2:
                cap = cv2.VideoCapture(str(path))
                if cap.isOpened():
                    w, h = cap.get(cv2.CAP_PROP_FRAME_WIDTH), cap.get(cv2.CAP_PROP_FRAME_HEIGHT)
                    cap.release()
                    if w > 0 and h > 0: return w/h
            else:
                with Image.open(path) as img: return img.size[0] / img.size[1]
        except: pass
        return 16/9
    
    def fit_media_in_placeholder(self, slide, placeholder, media_path, is_video=False, poster_image=None):
        """Place media into placeholder while preserving aspect ratio"""
        p_left, p_top, p_width, p_height = placeholder.left, placeholder.top, placeholder.width, placeholder.height
        aspect_ratio = self.get_aspect_ratio(Path(media_path), is_video)
        
        # Calculate size maintaining aspect ratio within placeholder bounds
        if aspect_ratio > p_width / p_height:
            scaled_w, scaled_h = p_width, p_width / aspect_ratio
        else:
            scaled_h, scaled_w = p_height, p_height * aspect_ratio
        
        # Center media within placeholder
        final_left = p_left + (p_width - scaled_w) / 2
        final_top = p_top + (p_height - scaled_h) / 2
        
        # Remove placeholder and add media
        placeholder._element.getparent().remove(placeholder._element)
        try:
            if is_video:
                if poster_image: slide.shapes.add_movie(str(media_path), final_left, final_top, scaled_w, scaled_h, poster_frame_image=str(poster_image))
                else: slide.shapes.add_movie(str(media_path), final_left, final_top, scaled_w, scaled_h)
            else:
                slide.shapes.add_picture(str(media_path), final_left, final_top, scaled_w, scaled_h)
        except:
            if is_video: slide.shapes.add_movie(str(media_path), final_left, final_top, scaled_w, scaled_h)
            else: slide.shapes.add_picture(str(media_path), final_left, final_top, scaled_w, scaled_h)
    
    def add_error_to_placeholder(self, slide, placeholder, error_msg):
        """Add error message to placeholder"""
        p_left, p_top, p_width, p_height = placeholder.left, placeholder.top, placeholder.width, placeholder.height
        placeholder._element.getparent().remove(placeholder._element)
        
        error_box = slide.shapes.add_textbox(p_left, p_top, p_width, p_height)
        error_box.text_frame.text = f"❌ EFFECT FAILED\n\n{error_msg}"
        for p in error_box.text_frame.paragraphs:
            p.font.size, p.alignment, p.font.color.rgb = Inches(16/72), PP_ALIGN.CENTER, RGBColor(255, 0, 0)
        error_box.fill.solid()
        error_box.fill.fore_color.rgb, error_box.line.color.rgb = RGBColor(255, 240, 240), RGBColor(255, 0, 0)
        error_box.line.width = Inches(0.02)
    
    def extract_video_key(self, filename: str, effect_name: str) -> str:
        stem = Path(filename).stem
        effect_clean = effect_name.lower().replace(' ', '_')
        pattern = re.escape(effect_clean) + r'_effect$'
        key = re.sub(pattern, '', stem, flags=re.IGNORECASE)
        for pattern in [r'_effect$', r'_generated$', r'_output$', r'_result$']:
            key = re.sub(pattern, '', key, flags=re.IGNORECASE)
        return normalize_key(key)
    
    def process_folders(self):
        base_folder = Path(self.config.get('base_folder', ''))
        if not base_folder.exists(): raise FileNotFoundError(f"Base folder not found: {base_folder}")
        
        pairs = []
        for task in self.config.get('tasks', []):
            effect, category = task.get('effect', ''), task.get('category', 'Unknown')
            folders = {k: base_folder / effect / v for k, v in 
                      {'src': 'Source', 'vid': 'Generated_Video', 'meta': 'Metadata'}.items()}
            
            if not folders['src'].exists(): continue
            logger.info(f"Processing: {effect}")
            
            # Get normalized file mappings
            images = {normalize_key(f.stem): f for f in folders['src'].iterdir() 
                     if f.suffix.lower() in {'.jpg', '.jpeg', '.png', '.webp'}}
            videos = {}
            if folders['vid'].exists():
                for f in folders['vid'].iterdir():
                    if f.suffix.lower() in {'.mp4', '.mov', '.avi'}:
                        key = self.extract_video_key(f.name, effect)
                        videos[key] = f
            
            logger.info(f"📊 Images: {len(images)}, Videos: {len(videos)}")
            
            # Create pairs
            for key, img in images.items():
                meta_file = folders['meta'] / f"{img.stem}_metadata.json"
                metadata = {}
                if meta_file.exists():
                    try:
                        with open(meta_file, 'r') as f: metadata = json.load(f)
                    except: pass
                
                vid = videos.get(key)
                if vid: logger.info(f"  ✓ Matched: {key}")
                else: logger.warning(f"  ❌ No video: {key}")
                
                pairs.append(MediaPair(
                    image_file=img.name, image_path=img, effect_name=effect, category=category,
                    video_file=vid.name if vid else None, video_path=vid, metadata=metadata,
                    failed=not vid or not metadata.get('success', False)
                ))
        
        return pairs
    
    def create_slide(self, ppt, pair, idx):
        # Try template first
        template_loaded = False
        try:
            if len(ppt.slides) >= 4:
                slide = ppt.slides.add_slide(ppt.slides[3].slide_layout)
                template_loaded = True
            else:
                slide = ppt.slides.add_slide(ppt.slide_layouts[6])
        except:
            slide = ppt.slides.add_slide(ppt.slide_layouts[6])
        
        if template_loaded:
            # Template mode: Use placeholders with aspect ratio preservation
            placeholders = sorted([p for p in slide.placeholders], key=lambda x: x.left if hasattr(x, 'left') else 0)
            
            # Update title
            for p in slide.placeholders:
                if p.placeholder_format.type == 1:
                    title = f"Effect {idx}: {pair.effect_name} - {pair.image_file}"
                    if pair.failed: title += " ❌ FAILED"
                    p.text = title
                    if pair.failed and p.text_frame.paragraphs:
                        p.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 0, 0)
                    break
            
            # Handle content placeholders
            content_placeholders = [p for p in placeholders if p.placeholder_format.type in {6, 7, 8, 13, 18, 19}]
            if len(content_placeholders) >= 2:
                self.fit_media_in_placeholder(slide, content_placeholders[0], str(pair.image_path), is_video=False)
                if pair.video_path and pair.video_path.exists():
                    self.fit_media_in_placeholder(slide, content_placeholders[1], str(pair.video_path), 
                                                is_video=True, poster_image=str(pair.image_path))
                else:
                    error_msg = pair.metadata.get('error_message', 'Effect failed') if pair.metadata else 'No processing'
                    self.add_error_to_placeholder(slide, content_placeholders[1], error_msg)
        
        else:
            # Fallback mode: Manual positioning
            title = f"Effect {idx}: {pair.effect_name} - {pair.image_file}" + (" ❌ FAILED" if pair.failed else "")
            title_box = slide.shapes.add_textbox(Cm(1), Cm(0.5), Cm(30), Cm(2))
            title_box.text_frame.text = title
            title_box.text_frame.paragraphs[0].font.size = Inches(18/72)
            if pair.failed: title_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 0, 0)
            
            # Manual positioning with aspect ratio
            img_x, img_y, img_w, img_h = self.calc_position(pair.image_path, *self.positions['img'])
            slide.shapes.add_picture(str(pair.image_path), img_x, img_y, img_w, img_h)
            
            if pair.video_path and pair.video_path.exists():
                vid_x, vid_y, vid_w, vid_h = self.calc_position(pair.video_path, *self.positions['vid'], True)
                try: slide.shapes.add_movie(str(pair.video_path), vid_x, vid_y, vid_w, vid_h, poster_frame_image=str(pair.image_path))
                except: slide.shapes.add_movie(str(pair.video_path), vid_x, vid_y, vid_w, vid_h)
            else:
                error_box = slide.shapes.add_textbox(Cm(self.positions['vid'][0]), Cm(self.positions['vid'][1]), 
                                                   Cm(self.positions['vid'][2]), Cm(self.positions['vid'][3]))
                error_msg = pair.metadata.get('error_message', 'Effect failed') if pair.metadata else 'No processing'
                error_box.text_frame.text = f"❌ EFFECT FAILED\n\n{error_msg}"
                for p in error_box.text_frame.paragraphs:
                    p.font.size, p.alignment, p.font.color.rgb = Inches(14/72), PP_ALIGN.CENTER, RGBColor(255, 0, 0)
                error_box.fill.solid()
                error_box.fill.fore_color.rgb, error_box.line.color.rgb = RGBColor(255, 240, 240), RGBColor(255, 0, 0)
        
        # **UPDATED METADATA** - 3-line format, smaller box
        processing_time = pair.metadata.get('processing_time_seconds', 'N/A') if pair.metadata else 'N/A'
        status = 'FAILED' if pair.failed else 'SUCCESS'
        meta_text = f"{pair.image_file}\n{pair.category} | {pair.effect_name}\n{processing_time}s \\ {status}"
        
        meta_box = slide.shapes.add_textbox(Cm(2), Cm(16.5), Cm(8), Cm(2))  # Smaller box
        meta_box.text_frame.text = meta_text
        meta_box.text_frame.word_wrap = True
        for p in meta_box.text_frame.paragraphs: 
            p.font.size, p.alignment = Inches(9/72), PP_ALIGN.LEFT
    
    def calc_position(self, path, x, y, w, h, is_video=False):
        """Fallback positioning for non-template mode"""
        try:
            ar = self.get_aspect_ratio(path, is_video)
            box_w, box_h = Cm(w), Cm(h)
            if ar > box_w / box_h:
                scaled_w, scaled_h = box_w, box_w / ar
            else:
                scaled_h, scaled_w = box_h, box_h * ar
            return Cm(x) + (box_w - scaled_w) / 2, Cm(y) + (box_h - scaled_h) / 2, scaled_w, scaled_h
        except:
            return Cm(x + 0.5), Cm(y + 0.5), Cm(w - 1), Cm(h - 1)
    
    def run(self):
        try:
            pairs = self.process_folders()
            if not pairs: return False
            
            # Create presentation
            try: ppt = Presentation(self.config.get('template_path', 'I2V templates.pptx'))
            except: ppt = Presentation()
            ppt.slide_width, ppt.slide_height = Cm(33.87), Cm(19.05)
            
            # Title slide
            base_name = Path(self.config.get('base_folder', '')).name
            match = re.match(r'(\d{4})\s+(.+)', base_name)
            date, project = match.groups() if match else (datetime.now().strftime("%m%d"), base_name)
            
            if ppt.slides and ppt.slides[0].shapes:
                ppt.slides[0].shapes[0].text_frame.text = f"[{date}] Vidu Effects\n{project}"
            
            # Content slides
            for i, pair in enumerate(pairs, 1): self.create_slide(ppt, pair, i)
            
            # Save with standardized filename
            filename = get_report_filename(base_name, "Vidu Effects")
            output_path = Path(self.config.get('output_directory', './')) / f"{filename}.pptx"
            ppt.save(str(output_path))
            logger.info(f"✓ Saved: {output_path} ({len(pairs)} effects)")
            return True
            
        except Exception as e:
            logger.error(f"Report failed: {e}")
            return False

if __name__ == "__main__":
    sys.exit(0 if ViduReportGenerator().run() else 1)
