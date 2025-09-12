import re, sys, json, logging
from pathlib import Path
from dataclasses import dataclass
from datetime import datetime
from concurrent.futures import ThreadPoolExecutor
from PIL import Image
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm, Inches, Pt
from pptx.dml.color import RGBColor

try: import cv2
except ImportError: cv2 = None

@dataclass
class MediaPair:
    image_file: str
    image_path: Path
    video_path: Path = None
    ref_video_path: Path = None
    source_video_path: Path = None
    metadata: dict = None
    ref_metadata: dict = None
    failed: bool = False
    ref_failed: bool = False

def get_runway_filename(folder_name, model_name="Runway", use_comparison=False, ref_folder_name=""):
    def extract_date_and_style(name):
        m = re.match(r'(?:\[(\d{4})\]|\[?(\d{4})\]?)\s*(.+)', name.strip())
        return (m.group(1) or m.group(2), m.group(3).strip()) if m else (datetime.now().strftime("%m%d"), name.strip())
    
    date, style = extract_date_and_style(folder_name)
    parts = [f"[{date}]"]
    if model_name.lower() not in style.lower(): parts.append(model_name)
    if use_comparison and ref_folder_name:
        _, ref_style = extract_date_and_style(ref_folder_name)
        parts.append(f"{style} vs {ref_style}")
    else:
        parts.append(style)
    return ' '.join(parts)

class RunwaySlideGenerator:
    def __init__(self, config_file="batch_runway_config.json"):
        self._ar_cache = {}
        logging.basicConfig(level=logging.INFO, format='%(message)s')
        self.logger = logging.getLogger(__name__)
        with open(config_file, 'r', encoding='utf-8') as f:
            self.config = json.load(f)
        self.logger.info(f"✓ Configuration loaded from {config_file}")

    def get_aspect_ratio(self, path, is_video=False):
        fn = path.name.lower()
        if '9_16' in fn or 'portrait' in fn: return 9/16
        if '1_1' in fn or 'square' in fn: return 1
        if '16_9' in fn or 'landscape' in fn: return 16/9
        key = str(path)
        if key in self._ar_cache: return self._ar_cache[key]
        try:
            if is_video and cv2:
                cap = cv2.VideoCapture(str(path))
                if cap.isOpened():
                    w, h = cap.get(3), cap.get(4)
                    cap.release()
                    if w > 0 and h > 0:
                        self._ar_cache[key] = w/h
                        return w/h
            else:
                with Image.open(path) as img:
                    ar = img.width / img.height
                    self._ar_cache[key] = ar
                    return ar
        except: pass
        return 16/9

    def add_media_to_slide(self, slide, ph, media=None, is_video=False, poster=None, error=None):
        l,t,w,h = ph.left, ph.top, ph.width, ph.height
        ph._element.getparent().remove(ph._element)
        if media and Path(media).exists():
            try:
                ar = self.get_aspect_ratio(Path(media), is_video)
                sw, sh = (w, w/ar) if ar > w/h else (h*ar, h)
                fl, ft = l+(w-sw)/2, t+(h-sh)/2
                if is_video:
                    if poster and Path(poster).exists():
                        slide.shapes.add_movie(media, fl, ft, sw, sh, poster_frame_image=poster)
                    else:
                        slide.shapes.add_movie(media, fl, ft, sw, sh)
                else:
                    slide.shapes.add_picture(media, fl, ft, sw, sh)
            except Exception as e:
                self._add_error_box(slide, l, t, w, h, f"Failed to load media: {e}")
        else:
            self._add_error_box(slide, l, t, w, h, error or "Media file not found")

    def _add_error_box(self, slide, l, t, w, h, msg):
        box = slide.shapes.add_textbox(l, t, w, h)
        box.text_frame.text = f"❌ GENERATION FAILED\n\n{msg}"
        for p in box.text_frame.paragraphs:
            p.font.size, p.alignment, p.font.color.rgb = Inches(16/72), PP_ALIGN.CENTER, RGBColor(255,0,0)
        box.fill.solid()
        box.fill.fore_color.rgb, box.line.color.rgb, box.line.width = RGBColor(255,240,240), RGBColor(255,0,0), Inches(0.02)

    def find_matching_video(self, base, videos):
        return videos.get(base) or next((p for n,p in videos.items() if n.startswith(base+'_')), None)

    def process_runway_batch(self, task):
        fp = Path(task['folder'])
        ref_fp = Path(task.get('reference_folder', '')) if task.get('reference_folder') else None
        use_cmp = task.get('use_comparison_template', False)
        
        folders = {k: fp/v for k,v in {'reference':'Reference','source':'Source','generated':'Generated_Video','metadata':'Metadata'}.items()}
        ref_folders = {k: ref_fp/v for k,v in {'video':'Generated_Video','metadata':'Metadata'}.items()} if ref_fp and use_cmp else {}
        
        if not folders['metadata'].exists(): return fp, [], task

        def load_files(folder, exts):
            return {f.stem: f for f in folder.iterdir() if f.suffix.lower() in exts} if folder.exists() else {}

        reference_images = load_files(folders['reference'], {'.jpg','.jpeg','.png','.webp'})
        source_videos = load_files(folders['source'], {'.mp4','.mov'})
        generated_videos = load_files(folders['generated'], {'.mp4','.mov'})
        metadata_files = load_files(folders['metadata'], {'.json'})
        
        ref_videos = load_files(ref_folders.get('video', Path()), {'.mp4','.mov'})
        ref_metadata = load_files(ref_folders.get('metadata', Path()), {'.json'})

        def load_meta(p):
            try:
                with open(p, 'r', encoding='utf-8') as f: return json.load(f)
            except: return {}

        pairs = []
        for stem, meta_path in metadata_files.items():
            md = load_meta(meta_path)
            if not md: continue
            
            ref_img_path = next((p for p in reference_images.values() if p.name == md.get('reference_image','')), None)
            src_vid_path = next((p for p in source_videos.values() if p.name == md.get('source_video','')), None)
            gen_vid_path = next((p for p in generated_videos.values() if p.name == md.get('generated_video','')), None)

            if ref_img_path:
                ref_vid_path, ref_md = None, {}
                if use_cmp:
                    rbase = stem.replace('_runway_metadata', '')
                    ref_vid_path = self.find_matching_video(rbase, ref_videos)
                    if rbase in ref_metadata: ref_md = load_meta(ref_metadata[rbase])

                pairs.append(MediaPair(ref_img_path.name, ref_img_path, gen_vid_path, ref_vid_path, src_vid_path, md, ref_md,
                           not gen_vid_path or not md.get('success', False), use_cmp and (not ref_vid_path or not ref_md.get('success', False))))
        return fp, pairs, task

    def create_runway_slide(self, ppt, pair, num, loaded, use_cmp=False):
        if loaded and len(ppt.slides) >= 4:
            slide = ppt.slides.add_slide(ppt.slides[3].slide_layout)
            for p in slide.placeholders:
                if p.placeholder_format.type == 1:
                    title = f"Generation {num}: {pair.image_file}" + (" ❌ FAILED" if pair.failed else "") + (" (REF ❌)" if use_cmp and pair.ref_failed else "")
                    p.text = title
                    if (pair.failed or (use_cmp and pair.ref_failed)) and p.text_frame.paragraphs:
                        p.text_frame.paragraphs[0].font.color.rgb = RGBColor(255,0,0)
                    break
            
            phs = sorted([p for p in slide.placeholders if p.placeholder_format.type in {6,7,8,13,18,19}], key=lambda x: getattr(x,'left',0))
            if len(phs) >= 3:
                self.add_media_to_slide(slide, phs[0], str(pair.image_path))
                self.add_media_to_slide(slide, phs[1], str(pair.source_video_path) if pair.source_video_path and pair.source_video_path.exists() else None, True, str(pair.image_path), "Source video not found")
                self.add_media_to_slide(slide, phs[2], str(pair.video_path) if pair.video_path and pair.video_path.exists() else None, True, str(pair.image_path), pair.metadata.get('error','Face swap generation failed') if pair.metadata else 'Generation failed')
            elif len(phs) >= 2:
                self.add_media_to_slide(slide, phs[0], str(pair.image_path))
                self.add_media_to_slide(slide, phs[1], str(pair.video_path) if pair.video_path and pair.video_path.exists() else None, True, str(pair.image_path), pair.metadata.get('error','Face swap generation failed') if pair.metadata else 'Generation failed')
            
            self._add_runway_metadata(slide, pair, use_cmp, len(phs) >= 3)
        else:
            slide = ppt.slides.add_slide(ppt.slide_layouts[6])
            title = f"Generation {num}: {pair.image_file}" + (" ❌ FAILED" if pair.failed else "") + (" (REF ❌)" if use_cmp and pair.ref_failed else "")
            tb = slide.shapes.add_textbox(Cm(2), Cm(1), Cm(20), Cm(2))
            tb.text_frame.text = title
            tb.text_frame.paragraphs[0].font.size = Inches(20/72)
            if pair.failed or (use_cmp and pair.ref_failed):
                tb.text_frame.paragraphs[0].font.color.rgb = RGBColor(255,0,0)
            
            pos = [(2.59,3.26), (13,3.26), (23.41,3.26)]
            w, h = 10, 10
            slide.shapes.add_picture(str(pair.image_path), Cm(pos[0][0]), Cm(pos[0][1]), Cm(w), Cm(h))
            
            for i, (vid_path, err_msg) in enumerate([(pair.source_video_path, "Source video not found"), 
                                                   (pair.video_path, pair.metadata.get('error','Face swap generation failed') if pair.metadata else 'Generation failed')], 1):
                if vid_path and vid_path.exists():
                    try:
                        slide.shapes.add_movie(str(vid_path), Cm(pos[i][0]), Cm(pos[i][1]), Cm(w), Cm(h), poster_frame_image=str(pair.image_path))
                    except:
                        slide.shapes.add_movie(str(vid_path), Cm(pos[i][0]), Cm(pos[i][1]), Cm(w), Cm(h))
                else:
                    self._add_error_box(slide, Cm(pos[i][0]), Cm(pos[i][1]), Cm(w), Cm(h), err_msg)
            
            self._add_runway_metadata(slide, pair, use_cmp, True)

    def _add_runway_metadata(self, slide, pair, use_cmp, is_three_col):
        meta = [f"{k}: {pair.metadata.get(k,'N/A') if k != 'success' else ('✓' if pair.metadata.get(k) else '❌')}" for k in ['prompt','reference_image','source_video','model','processing_time_seconds','success']] if pair.metadata else ["No metadata available"]
        if pair.metadata and 'processing_time_seconds' in meta[-2]: meta[-2] += 's'
        if use_cmp and pair.ref_metadata: meta.append(f"Ref Time: {pair.ref_metadata.get('processing_time_seconds','N/A')}s")
        
        pos = (2.32, 15.24) if is_three_col else (5.19, 15.99)
        box = slide.shapes.add_textbox(Cm(pos[0]), Cm(pos[1]), Cm(7.29), Cm(3.06))
        box.text_frame.text = "\n".join(meta)
        for p in box.text_frame.paragraphs: p.font.size = Inches(10/72)

    def add_links_to_presentation(self, ppt, task):
        if not ppt.slides: return
        info_box = next((s for s in ppt.slides[0].shapes if hasattr(s,'text_frame') and s.text_frame.text and any(k in s.text_frame.text.lower() for k in ['design','testbed','source'])), None)
        if not info_box: info_box = ppt.slides[0].shapes.add_textbox(Cm(5), Cm(13), Cm(20), Cm(4))
        info_box.text_frame.clear()
        
        links = [("Design: ","Link",task.get('design_link','')), ("Testbed: ",self.config.get('testbed','http://192.168.4.3:8000/video_effect/'),self.config.get('testbed','http://192.168.4.3:8000/video_effect/')), ("Source + Video: ","Link",task.get('source_video_link',''))]
        for i, (pre, txt, url) in enumerate(links):
            para = info_box.text_frame.paragraphs[0] if i == 0 else info_box.text_frame.add_paragraph()
            if url:
                para.clear()
                r1, r2 = para.add_run(), para.add_run()
                r1.text, r1.font.size = pre, Inches(24/72)
                r2.text, r2.font.size = url if "Testbed:" in pre else txt, Inches(24/72)
                r2.hyperlink.address = url
                para.alignment = PP_ALIGN.CENTER
            else:
                para.text, para.font.size, para.alignment = f"{pre}{txt}", Inches(24/72), PP_ALIGN.CENTER

    def create_presentation(self, folder_path, media_pairs, task):
        if not media_pairs: return False
        use_cmp = task.get('use_comparison_template', False)
        ref_folder = task.get('reference_folder', '')
        template_path = self.config.get('comparison_template_path' if use_cmp else 'template_path', 'I2V Comparison Template.pptx' if use_cmp else 'I2V templates.pptx')
        
        try:
            ppt = Presentation(template_path) if Path(template_path).exists() else Presentation()
            loaded = Path(template_path).exists()
        except:
            ppt, loaded = Presentation(), False
        
        ppt.slide_width, ppt.slide_height = Cm(33.87), Cm(19.05)
        
        if ppt.slides:
            def extract_date_style(n):
                m = re.match(r'(?:\[(\d{4})\]|(\d{4}))?\s*(.*)', n.strip())
                d, s = (m.group(1) or m.group(2), m.group(3).strip()) if m else (None, '')
                return (d or datetime.now().strftime("%m%d"), s or n.strip())
            
            folder_name = folder_path.name
            slide = ppt.slides[0]
            for s in [s for s in slide.shapes if hasattr(s,'text_frame') and (not s.text_frame.text or "Results" in s.text_frame.text or any(k in s.text_frame.text for k in ["Runway","vs","["]))]:
                slide.shapes._spTree.remove(s._element)
            
            if use_cmp and ref_folder:
                d, s1 = extract_date_style(folder_name)
                _, s2 = extract_date_style(Path(ref_folder).name)
                tb = slide.shapes.add_textbox(Cm(5), Cm(2), Cm(24), Cm(8))
                tf = tb.text_frame
                p1, p2 = tf.paragraphs[0], tf.add_paragraph()
                p1.text, p2.text = f"[{d}] Runway", f"{s1} vs {s2}"
                for p, size in [(p1, 60), (p2, 40)]:
                    p.font.size, p.font.bold, p.alignment = Pt(size), True, PP_ALIGN.CENTER
            else:
                d, s = extract_date_style(folder_name)
                tb = slide.shapes.add_textbox(Cm(5), Cm(3), Cm(24), Cm(8))
                tf = tb.text_frame
                p1, p2 = tf.paragraphs[0], tf.add_paragraph()
                p1.text, p2.text = f"[{d}] Runway", s
                for p, size in [(p1, 60), (p2, 40)]:
                    p.font.size, p.font.bold, p.alignment = Pt(size), True, PP_ALIGN.CENTER
        
        self.add_links_to_presentation(ppt, task)
        for i, pair in enumerate(media_pairs, 1):
            self.create_runway_slide(ppt, pair, i, loaded, use_cmp)
        
        filename = get_runway_filename(folder_path.name, "Runway", use_cmp, Path(ref_folder).name if use_cmp and ref_folder else "")
        opath = Path(self.config.get('output_directory', './')) / f"{filename}.pptx"
        ppt.save(str(opath))
        self.logger.info(f"✓ Saved: {opath}")
        return True

    def validate_task_structure(self, task):
        folder = Path(task['folder'])
        required = ['Reference','Source','Generated_Video','Metadata']
        missing = [f for f in required if not (folder/f).exists()]
        if missing: return False, f"Missing folders: {', '.join(missing)}"
        
        counts = []
        for name, exts in [('face references', ['*.jpg','*.jpeg','*.png','*.webp']), ('source videos', ['*.mp4','*.mov']), ('metadata files', ['*.json'])]:
            files = []
            folder_path = folder / (['Reference','Source','Metadata'][len(counts)])
            for ext in exts: files.extend(list(folder_path.glob(ext)))
            if not files: return False, f"No {name} found in {folder_path.name} folder"
            counts.append(len(files))
        
        return True, f"Ready: {counts[0]} face references, {counts[1]} source videos, {counts[2]} metadata files"

    def process_task(self, task, i, total):
        folder = Path(task['folder'])
        self.logger.info(f"🎯 Task {i}/{total}: {folder.name}")
        try:
            valid, msg = self.validate_task_structure(task)
            if not valid:
                self.logger.error(f"❌ Task {i}: {msg}")
                return False
            
            fp, pairs, proc_task = self.process_runway_batch(task)
            if not pairs:
                self.logger.warning(f"⚠️ Task {i}: No media pairs found")
                return False
            
            success = self.create_presentation(fp, pairs, proc_task)
            self.logger.info(f"{'✅' if success else '❌'} Task {i}: {'Generated face swap slides with '+str(len(pairs))+' pairs' if success else 'Slide generation failed'}")
            return success
        except Exception as e:
            self.logger.error(f"❌ Task {i} error: {e}")
            return False

    def generate_all_presentations(self, parallel=True):
        tasks = self.config.get('tasks', [])
        if not tasks:
            self.logger.warning("⚠️ No tasks found in configuration")
            return 0
        
        self.logger.info(f"🚀 Generating face swap slides for {len(tasks)} tasks")
        if parallel and len(tasks) > 1:
            with ThreadPoolExecutor(max_workers=min(4, len(tasks))) as ex:
                success_count = sum(1 for f in [ex.submit(self.process_task, t, i, len(tasks)) for i, t in enumerate(tasks, 1)] if f.result())
        else:
            success_count = sum(1 for i, t in enumerate(tasks, 1) if self.process_task(t, i, len(tasks)))
        
        self.logger.info(f"🎉 Completed: {success_count}/{len(tasks)} face swap presentations generated")
        if success_count > 0:
            self.logger.info(f"📁 Saved to: {Path(self.config.get('output_directory', './')).absolute()}")
        return success_count

    def run(self, parallel=True):
        self.logger.info("🎬 Starting Runway Face Swap Slide Generator")
        try:
            return self.generate_all_presentations(parallel) > 0
        except Exception as e:
            self.logger.error(f"❌ Critical error: {e}")
            return False

def main():
    import argparse
    parser = argparse.ArgumentParser(description='Generate slides from Runway face swap batch processing results')
    parser.add_argument('--config', '-c', default='batch_runway_config.json', help='Config file (default: batch_runway_config.json)')
    parser.add_argument('--sequential', '-s', action='store_true', help='Process tasks sequentially instead of parallel')
    args = parser.parse_args()
    try:
        gen = RunwaySlideGenerator(args.config)
        sys.exit(0 if gen.run(parallel=not args.sequential) else 1)
    except Exception as e:
        print(f"❌ Fatal error: {e}")
        sys.exit(1)

if __name__ == '__main__': main()
