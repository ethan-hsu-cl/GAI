"""PowerPoint presentation building service."""
import json
import tempfile
from pathlib import Path
from PIL import Image
from pptx import Presentation
from pptx.util import Cm, Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from typing import List, Dict, Optional
import logging

try:
    import cv2
except ImportError:
    cv2 = None

class PresentationBuilder:
    """Builds PowerPoint presentations from media pairs."""
    
    def __init__(self, template_path: str):
        self.template_path = template_path
        self.logger = logging.getLogger(__name__)
    
    def create_presentation(self, media_pairs: List, output_path: Path) -> bool:
        """Create presentation from media pairs."""
        try:
            # Load template
            prs = Presentation(self.template_path)
            
            # Process each media pair
            for pair in media_pairs:
                self._add_slide_for_pair(prs, pair)
            
            # Save presentation
            prs.save(output_path)
            return True
            
        except Exception as e:
            self.logger.error(f"Presentation creation failed: {e}")
            return False
    
    def _add_slide_for_pair(self, prs: Presentation, pair) -> None:
        """Add slide for a media pair."""
        # Implementation would depend on the specific slide layout
        # This is a simplified version
        slide_layout = prs.slide_layouts[1]  # Assuming layout 1
        slide = prs.slides.add_slide(slide_layout)
        
        # Add title
        title = slide.shapes.title
        title.text = f"Processing Results: {pair.source_file}"
        
        # Add source image/video
        if pair.source_path.exists():
            self._add_media_to_slide(slide, pair.source_path, 
                                   left=Cm(1), top=Cm(3), width=Cm(8))
        
        # Add generated content
        for i, gen_path in enumerate(pair.generated_paths[:3]):  # Limit to 3
            if gen_path.exists():
                self._add_media_to_slide(slide, gen_path,
                                       left=Cm(10 + i*5), top=Cm(3), width=Cm(4))
    
    def _add_media_to_slide(self, slide, media_path: Path, 
                           left: Cm, top: Cm, width: Cm) -> None:
        """Add media (image or video) to slide."""
        try:
            if media_path.suffix.lower() in ['.jpg', '.jpeg', '.png', '.bmp']:
                # Add image
                slide.shapes.add_picture(str(media_path), left, top, width=width)
            elif media_path.suffix.lower() in ['.mp4', '.mov', '.avi']:
                # For videos, add a thumbnail or placeholder
                thumbnail_path = self._create_video_thumbnail(media_path)
                if thumbnail_path:
                    slide.shapes.add_picture(str(thumbnail_path), left, top, width=width)
        except Exception as e:
            self.logger.warning(f"Could not add media {media_path}: {e}")
    
    def _create_video_thumbnail(self, video_path: Path) -> Optional[Path]:
        """Create thumbnail from video."""
        if not cv2:
            return None
            
        try:
            cap = cv2.VideoCapture(str(video_path))
            ret, frame = cap.read()
            if ret:
                thumbnail_path = Path(tempfile.mkdtemp()) / f"{video_path.stem}_thumb.jpg"
                cv2.imwrite(str(thumbnail_path), frame)
                cap.release()
                return thumbnail_path
            cap.release()
        except Exception:
            pass
        return None
