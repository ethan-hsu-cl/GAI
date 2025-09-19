"""Report generator with EXACT workflow from unified_report_generator.py"""
import json
import logging
import os
import re
import tempfile
from datetime import datetime
from pathlib import Path
from typing import Dict, List, Optional

try:
    import cv2
except ImportError:
    cv2 = None
    
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Cm, Inches, Pt
from pptx.enum.text import PP_ALIGN

from .config_manager import ConfigManager
from ..models.media_pair import MediaPair

class ReportGenerator:
    """Report generator with EXACT workflow from unified_report_generator.py"""
    
    def __init__(self, api_name: str, config_file: str = None):
        self.api_name = api_name
        self.config_file = config_file
        self.config = None
        self.report_definitions = {}
        self.frame_cache = {}
        self.ar_cache = {}
        self.logger = logging.getLogger(__name__)
        
        self._load_config()
        self._load_report_definitions()
    
    def _load_config(self):
        """Load configuration from file."""
        if self.config_file:
            self.config = ConfigManager.load_config(self.config_file)
    
    def _load_report_definitions(self):
        """Load report definitions from api_definitions.json"""
        definition_paths = ['core/api_definitions.json', 'api_definitions.json']
        
        for def_path in definition_paths:
            try:
                with open(def_path, 'r', encoding='utf-8') as f:
                    all_definitions = json.load(f)
                    self.report_definitions = all_definitions.get(self.api_name, {}).get('report', {})
                    return
            except Exception:
                continue
        
        # Default report definitions
        self.report_definitions = {
            'enabled': True,
            'template_path': 'templates/I2V templates.pptx',
            'comparison_template_path': 'templates/I2V Comparison Template.pptx',
            'output_directory': '/Users/ethanhsu/Desktop/GAI/Report',
            'use_comparison': self.api_name in ['kling', 'nano_banana', 'runway']
        }
    
    def get_filename(self, folder, model):
        """Generate filename - FIXED to avoid duplicating date in folder name"""
        folder_name = Path(folder).name if hasattr(folder, 'name') else str(folder)
        
        # Extract date from folder name and remove it from the style name
        m = re.match(r'(\d{4})\s*(.*)', folder_name)  # Match date and capture the rest
        if m:
            date_digits = m.group(1)
            remaining_text = m.group(2).strip()  # Get text after the date
            
            # If it looks like a year, use current month/day
            if date_digits.startswith('20'):
                d = datetime.now().strftime('%m%d')
            else:
                d = date_digits
            
            # Use remaining text (without date) as style name
            s = remaining_text if remaining_text else folder_name
        else:
            # No date found, use current date and full folder name
            d = datetime.now().strftime('%m%d')
            s = folder_name
        
        parts = [f"[{d}]"]
        if model and model.lower() not in s.lower():
            parts.append(model)
        parts.append(s)
        return " ".join(parts)

    def get_cmp_filename(self, folder1: str, folder2: str, model: str):
        """Generate comparison filename - FIXED to avoid duplicating date"""
        folder1_name = Path(folder1).name if hasattr(folder1, 'name') else str(folder1)
        folder2_name = Path(folder2).name if hasattr(folder2, 'name') else str(folder2)
        
        # Extract date from first folder and remove it from the style name
        m1 = re.match(r'(\d{4})\s*(.*)', folder1_name)
        if m1:
            date_digits = m1.group(1)
            remaining_text = m1.group(2).strip()
            
            if date_digits.startswith('20'):
                d = datetime.now().strftime('%m%d')
            else:
                d = date_digits
            
            s1 = remaining_text if remaining_text else folder1_name
        else:
            d = datetime.now().strftime('%m%d')
            s1 = folder1_name
        
        # Extract style name from second folder (remove date if present)
        m2 = re.match(r'(\d{4})\s*(.*)', folder2_name)
        if m2:
            s2 = m2.group(2).strip() if m2.group(2).strip() else folder2_name
        else:
            s2 = folder2_name
        
        parts = [f"[{d}]"]
        if model and model.lower() not in s1.lower():
            parts.append(model)
        parts.append(f"{s1} vs {s2}")
        return " ".join(parts)

    
    def calc_pos(self, path: Path, pos: tuple) -> tuple:
        """Calculate positioned media - EXACT from unified_report_generator.py"""
        x, y, w, h = pos
        try:
            with Image.open(path) as img:
                ar = img.width / img.height
            bw, bh = Cm(w), Cm(h)
            sw, sh = (bw, bw/ar) if ar > bw/bh else (bh*ar, bh)
            return Cm(x) + (bw-sw)/2, Cm(y) + (bh-sh)/2, sw, sh
        except:
            return Cm(x) + Cm(0.5), Cm(y) + Cm(0.5), Cm(w) - Cm(1), Cm(h) - Cm(1)
    
    def get_aspect_ratio(self, path, is_video=False):
        """Aspect ratio calculation - EXACT from unified_report_generator.py"""
        fn = str(path).lower()
        if '9-16' in fn or 'portrait' in fn:
            return 9/16
        if '1-1' in fn or 'square' in fn:
            return 1
        if '16-9' in fn or 'landscape' in fn:
            return 16/9
        
        key = str(path)
        if key in self.ar_cache:
            return self.ar_cache[key]
        
        try:
            if is_video and cv2:
                cap = cv2.VideoCapture(str(path))
                if cap.isOpened():
                    w, h = cap.get(3), cap.get(4)
                    cap.release()
                    if w > 0 and h > 0:
                        self.ar_cache[key] = w/h
                        return w/h
            else:
                with Image.open(path) as img:
                    ar = img.width / img.height
                    self.ar_cache[key] = ar
                    return ar
        except:
            pass
        return 16/9
    
    def extract_first_frame(self, video_path):
        """Extract first frame - EXACT from unified_report_generator.py"""
        if not cv2:
            return None
            
        video_key = str(video_path)
        if video_key in self.frame_cache:
            return self.frame_cache[video_key]
        
        try:
            cap = cv2.VideoCapture(str(video_path))
            if not cap.isOpened():
                return None
            
            ret, frame = cap.read()
            cap.release()
            
            if not ret or frame is None:
                return None
            
            temp_dir = tempfile.gettempdir()
            frame_filename = f"frame_{video_path.stem}_{hash(video_key) % 10000}.jpg"
            frame_path = Path(temp_dir) / frame_filename
            
            cv2.imwrite(str(frame_path), frame)
            self.frame_cache[video_key] = str(frame_path)
            return str(frame_path)
            
        except Exception as e:
            self.logger.warning(f"Failed to extract frame from {video_path}: {e}")
            return None
    
    def create_title_slide(self, ppt: Presentation, task: Dict, use_comparison: bool):
        """Create title slide - FIXED to avoid duplicating date in title"""
        if not ppt.slides:
            return
        
        # Get folder name for title generation
        if self.api_name in ['vidu_effects', 'vidu_reference']:
            folder_name = Path(self.config.get('base_folder', '')).name
        else:
            folder_name = task.get('folder', Path(self.config.get('base_folder', '')).name)
            if isinstance(folder_name, str):
                folder_name = Path(folder_name).name
        
        # API display names
        api_display_names = {
            'kling': 'Kling 2.1',
            'nano_banana': 'Nano Banana',
            'runway': 'Runway',
            'vidu_effects': 'Vidu Effects',
            'vidu_reference': 'Vidu Reference',
            'genvideo': 'GenVideo'
        }
        api_display = api_display_names.get(self.api_name, self.api_name.title())
        
        # Generate title using FIXED patterns
        if use_comparison and task.get('reference_folder'):
            ref_name = Path(task['reference_folder']).name
            title = self.get_cmp_filename(folder_name, ref_name, api_display)
        else:
            title = self.get_filename(folder_name, api_display)
        
        slide = ppt.slides[0]
        
        # ✅ Update template placeholder directly
        if ppt.slides and ppt.slides[0].shapes:
            ppt.slides[0].shapes[0].text_frame.text = title
        else:
            # Fallback: Extract date and style for proper formatting
            m = re.match(r'(\d{4})\s*(.*)', str(folder_name))
            if m:
                d = m.group(1)
                if d.startswith('20'):
                    d = datetime.now().strftime('%m%d')
                s = m.group(2).strip() if m.group(2).strip() else str(folder_name)
            else:
                d = datetime.now().strftime('%m%d')
                s = str(folder_name)
            
            tb = slide.shapes.add_textbox(Cm(5), Cm(2), Cm(24), Cm(8))
            tf = tb.text_frame
            
            # First line: [Date] [API Name] - Font 60
            p1 = tf.paragraphs[0]
            p1.text = f"[{d}] {api_display}"
            p1.font.size = Pt(60)
            p1.font.bold = True
            p1.alignment = PP_ALIGN.CENTER
            
            # Second line: [Style Name] - Font 40
            p2 = tf.add_paragraph()
            p2.text = s
            p2.font.size = Pt(40)
            p2.font.bold = True
            p2.alignment = PP_ALIGN.CENTER
        
        # Add links section
        self.add_links_working_pattern(ppt, task)

    
    def add_links_working_pattern(self, ppt: Presentation, task: Dict):
        """Add hyperlinks - EXACT from unified_report_generator.py"""
        if not ppt.slides:
            return
        
        slide = ppt.slides[0]
        
        # ✅ CRITICAL FIX: Find existing infobox or create new one
        infobox = None
        for s in slide.shapes:
            if (hasattr(s, 'text_frame') and s.text_frame.text and 
                any(x in s.text_frame.text.lower() for x in ['design', 'testbed', 'source'])):
                infobox = s
                break
        
        if not infobox:
            infobox = slide.shapes.add_textbox(Cm(5), Cm(13), Cm(20), Cm(4))
            
        infobox.text_frame.clear()
        
        # Get links from task/config
        design_link = task.get('design_link', '')
        testbed_url = self.config.get('testbed', f'http://192.168.4.3:8000/{self.api_name}')
        source_link = task.get('source_video_link', '')
        
        lines = [
            "Design: Link",
            f"Testbed: {testbed_url}",
            "Source + Video: Link"
        ]
        
        for i, line in enumerate(lines):
            p = infobox.text_frame.paragraphs[0] if i == 0 else infobox.text_frame.add_paragraph()
            p.text, p.font.size, p.alignment = line, Inches(24/72), PP_ALIGN.CENTER
            
            if "Design" in line and task.get('design_link'):
                self.add_hyperlink(p, "Design: ", "Link", task['design_link'])
            elif "Testbed" in line:
                self.add_hyperlink(p, "Testbed: ", testbed_url, testbed_url)
            elif "Source Video" in line and task.get('source_video_link'):
                self.add_hyperlink(p, "Source + Video: ", "Link", task['source_video_link'])
    
    def add_hyperlink(self, para, pre, linktext, url):
        """Simple hyperlink creation - EXACT from unified_report_generator.py"""
        para.clear()
        r1 = para.add_run()
        r1.text, r1.font.size = pre, Inches(24/72)
        r2 = para.add_run()
        r2.text, r2.font.size = linktext, Inches(24/72)
        r2.hyperlink.address = url
        para.alignment = PP_ALIGN.CENTER
    
    def add_media_standard(self, slide, placeholder, media_path, is_video=False):
        """Add media to placeholder - EXACT from unified_report_generator.py"""
        p_left, p_top, p_width, p_height = placeholder.left, placeholder.top, placeholder.width, placeholder.height
        placeholder.element.getparent().remove(placeholder.element)
        
        try:
            aspect_ratio = self.get_aspect_ratio(Path(media_path), is_video)
            if aspect_ratio > p_width / p_height:
                scaled_w, scaled_h = p_width, p_width / aspect_ratio
            else:
                scaled_w, scaled_h = p_height * aspect_ratio, p_height
            
            final_left = p_left + (p_width - scaled_w) / 2
            final_top = p_top + (p_height - scaled_h) / 2
            
            if is_video:
                first_frame_path = self.extract_first_frame(Path(media_path))
                if first_frame_path and Path(first_frame_path).exists():
                    slide.shapes.add_movie(str(media_path), final_left, final_top, scaled_w, scaled_h, 
                                         poster_frame_image=first_frame_path)
                else:
                    slide.shapes.add_movie(str(media_path), final_left, final_top, scaled_w, scaled_h)
            else:
                slide.shapes.add_picture(str(media_path), final_left, final_top, scaled_w, scaled_h)
        except Exception as e:
            self.logger.warning(f"Failed to add media with aspect ratio, using fallback: {e}")
            if is_video:
                slide.shapes.add_movie(str(media_path), p_left + p_width*0.1, p_top + p_height*0.1, 
                                     p_width*0.8, p_height*0.8)
            else:
                slide.shapes.add_picture(str(media_path), p_left + p_width*0.1, p_top + p_height*0.1, 
                                       p_width*0.8, p_height*0.8)
    
    def add_error_standard(self, slide, placeholder, error_msg):
        """Add error to placeholder - EXACT from unified_report_generator.py"""
        p_left, p_top, p_width, p_height = placeholder.left, placeholder.top, placeholder.width, placeholder.height
        placeholder.element.getparent().remove(placeholder.element)
        
        error_box = slide.shapes.add_textbox(p_left, p_top, p_width, p_height)
        error_box.text_frame.text = f"GENERATION FAILED\n{error_msg}"
        for p in error_box.text_frame.paragraphs:
            p.font.size, p.alignment, p.font.color.rgb = Inches(16/72), PP_ALIGN.CENTER, RGBColor(255, 0, 0)
        error_box.fill.solid()
        error_box.fill.fore_color.rgb = RGBColor(255, 240, 240)
        error_box.line.color.rgb = RGBColor(255, 0, 0)
        error_box.line.width = Inches(0.02)
    
    def create_standard_slide(self, ppt, pair, index, template_loaded, use_comparison):
        """Create standard slide - EXACT from unified_report_generator.py"""
        if template_loaded and len(ppt.slides) >= 4:
            # ✅ Use 4th slide (index 3) as base
            slide = ppt.slides.add_slide(ppt.slides[3].slide_layout)
            
            # Handle title placeholder
            for p in slide.placeholders:
                if p.placeholder_format.type == 1:  # Title
                    title = f"Generation {index}: {pair.source_file}"
                    if pair.failed:
                        title += " - FAILED"
                    p.text = title
                    if pair.failed and p.text_frame.paragraphs:
                        p.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 0, 0)
                    break
            
            # Get media placeholders and sort by position
            phs = sorted([p for p in slide.placeholders 
                         if p.placeholder_format.type in [6, 7, 8, 13, 18, 19]], 
                        key=lambda x: getattr(x, 'left', 0))
            
            if use_comparison and len(phs) >= 3:
                # 3-way comparison: Source, Generated, Reference
                self.add_media_standard(slide, phs[0], str(pair.source_path), False)
                
                if pair.generated_paths and pair.generated_paths[0].exists():
                    is_video = str(pair.generated_paths[0]).lower().endswith(('.mp4', '.mov', '.avi'))
                    self.add_media_standard(slide, phs[1], str(pair.generated_paths[0]), is_video)
                else:
                    self.add_error_standard(slide, phs[1], "Generation failed")
                
                if pair.reference_paths and pair.reference_paths[0].exists():
                    self.add_media_standard(slide, phs[2], str(pair.reference_paths[0]), False)
                else:
                    self.add_error_standard(slide, phs[2], "No reference")
                    
            elif len(phs) >= 2:
                # 2-way comparison: Source, Generated
                self.add_media_standard(slide, phs[0], str(pair.source_path), False)
                
                if pair.generated_paths and pair.generated_paths[0].exists():
                    is_video = str(pair.generated_paths[0]).lower().endswith(('.mp4', '.mov', '.avi'))
                    self.add_media_standard(slide, phs[1], str(pair.generated_paths[0]), is_video)
                else:
                    self.add_error_standard(slide, phs[1], "Generation failed")
            
            # ✅ Add metadata box with EXACT positioning from original
            self.add_metadata_box_exact(slide, pair, len(phs))
        else:
            # Fallback to blank slide if template not loaded
            slide = ppt.slides.add_slide(ppt.slide_layouts[6])
            self.logger.warning("Using blank slide layout as fallback")
    
    def add_metadata_box_exact(self, slide, pair: MediaPair, media_count: int):
        """Add metadata box - EXACT positioning from unified_report_generator.py"""
        if self.api_name == 'nano_banana':
            # ✅ EXACT Nano Banana metadata positioning: Cm(5), Cm(16), Cm(12), Cm(3)
            mb = slide.shapes.add_textbox(Cm(5), Cm(16), Cm(12), Cm(3))
            response_id = pair.metadata.get('response_id', 'N/A') if pair.metadata else 'N/A'
            proc_time = pair.metadata.get('processing_time_seconds', 'N/A') if pair.metadata else 'N/A'
            mb.text_frame.text = f"File Name: {pair.source_file}\nID: {response_id}\n{proc_time}s"
            for p in mb.text_frame.paragraphs:
                p.font.size = Inches(10/72)
                
        elif self.api_name == 'kling':
            # ✅ EXACT Kling metadata positioning: Cm(2), Cm(16), Cm(15), Cm(3)
            metabox = slide.shapes.add_textbox(Cm(2), Cm(16), Cm(15), Cm(3))
            if pair.metadata:
                task_id = pair.metadata.get('task_id', 'N/A')
                proc_time = pair.metadata.get('processing_time_seconds', 'N/A')
                success = "✓" if pair.metadata.get('success', False) else "✗"
                model = pair.metadata.get('model', 'N/A')
                prompt = pair.metadata.get('prompt', 'N/A')
                
                meta_lines = [
                    f"Task: {task_id}",
                    f"Model: {model}",
                    f"Prompt: {prompt[:50]}..." if len(str(prompt)) > 50 else f"Prompt: {prompt}",
                    f"Time: {proc_time}s",
                    f"Status: {success}"
                ]
            else:
                meta_lines = ["No metadata available"]
            
            metabox.text_frame.text = "\n".join(meta_lines)
            metabox.text_frame.word_wrap = True
            for para in metabox.text_frame.paragraphs:
                para.font.size = Inches(10/72)
                
        elif self.api_name == 'runway':
            # ✅ EXACT Runway metadata positioning (varies by layout)
            pos = (2.32, 15.24) if media_count >= 3 else (5.19, 15.99)
            box = slide.shapes.add_textbox(Cm(pos[0]), Cm(pos[1]), Cm(7.29), Cm(3.06))
            
            meta = [f"{k}: {pair.metadata.get(k, 'N/A')}" if k != 'success' else 
                   ("✓" if pair.metadata.get(k) else "✗") 
                   for k in ['prompt', 'reference_image', 'source_video', 'model', 'processing_time_seconds', 'success']
                   if pair.metadata] or ["No metadata available"]
                   
            if pair.metadata and len(meta) > 1 and 'processing_time_seconds' in meta[-2]:
                meta[-2] = meta[-2].replace('processing_time_seconds', 'processing_time_seconds') + 's'
                
            box.text_frame.text = "\n".join(meta)
            for p in box.text_frame.paragraphs:
                p.font.size = Inches(10/72)
        else:
            # Generic metadata for other APIs
            meta_box = slide.shapes.add_textbox(Cm(18), Cm(16), Cm(13), Cm(3))
            meta_lines = [f"File: {pair.source_file}"]
            if pair.metadata:
                proc_time = pair.metadata.get('processing_time_seconds', 'N/A')
                success = "✓" if pair.metadata.get('success', False) else "✗"
                meta_lines.extend([f"Time: {proc_time}s", f"Status: {success}"])
                
            meta_box.text_frame.text = "\n".join(meta_lines)
            for para in meta_box.text_frame.paragraphs:
                para.font.size = Inches(10/72)
    
    def create_presentation(self, pairs: List[MediaPair], output_path: Path) -> bool:
        """Create presentation - EXACT workflow from unified_report_generator.py"""
        if not pairs:
            self.logger.warning("No media pairs to process")
            return False
        
        # Get task info from config for title page
        task = {}
        if self.config and self.config.get('tasks'):
            task = self.config['tasks'][0]  # Use first task for title info
        
        # ✅ Check config for comparison template selection
        use_comparison = task.get('use_comparison_template', False)
        template_key = 'comparison_template_path' if use_comparison else 'template_path'
        template_path = (self.config.get(template_key) if self.config else None) or (
            self.report_definitions.get(template_key, 
            'templates/I2V Comparison Template.pptx' if use_comparison else 'templates/I2V templates.pptx'))
        
        self.logger.info(f"Template selection: use_comparison={use_comparison}, template={template_path}")
        
        # ✅ Load template
        try:
            ppt = Presentation(template_path) if Path(template_path).exists() else Presentation()
            template_loaded = Path(template_path).exists()
            self.logger.info(f"Template loaded: {template_path}")
        except Exception as e:
            self.logger.warning(f"Template load failed: {e}, using blank presentation")
            ppt = Presentation()
            template_loaded = False
        
        # Set slide dimensions
        ppt.slide_width, ppt.slide_height = Cm(33.87), Cm(19.05)
        
        # ✅ Create title slide (update template, don't duplicate)
        self.create_title_slide(ppt, task, use_comparison)
        
        # ✅ Create content slides using 4th slide as base
        for i, pair in enumerate(pairs, 1):
            self.create_standard_slide(ppt, pair, i, template_loaded, use_comparison)
        
        # ✅ Save with correct filename pattern
        try:
            # Generate correct filename
            if self.api_name in ['vidu_effects', 'vidu_reference']:
                folder_name = Path(self.config.get('base_folder', '')).name
            else:
                first_task = self.config['tasks'][0] if self.config and self.config.get('tasks') else {}
                folder_name = first_task.get('folder', '')
                if isinstance(folder_name, str):
                    folder_name = Path(folder_name).name
            
            api_display_names = {
                'kling': 'Kling 2.1',
                'nano_banana': 'Nano Banana',
                'runway': 'Runway',
                'vidu_effects': 'Vidu Effects',
                'vidu_reference': 'Vidu Reference',
                'genvideo': 'GenVideo'
            }
            api_display = api_display_names.get(self.api_name, self.api_name.title())
            
            # Generate correct filename
            if use_comparison and task.get('reference_folder'):
                ref_name = Path(task['reference_folder']).name
                filename = self.get_cmp_filename(folder_name, ref_name, api_display)
            else:
                filename = self.get_filename(folder_name, api_display)
            
            # Get output directory from config
            output_directory = (self.config.get('output_directory') if self.config else None) or (
                self.report_definitions.get('output_directory', '/Users/ethanhsu/Desktop/GAI/Report'))
            
            final_output_path = Path(output_directory) / f"{filename}.pptx"
            final_output_path.parent.mkdir(parents=True, exist_ok=True)
            
            ppt.save(str(final_output_path))
            self.logger.info(f"Report saved: {final_output_path}")
            print(f"Report generated: {final_output_path}")
            return True
        except Exception as e:
            self.logger.error(f"Failed to save presentation: {e}")
            return False
    
    def cleanup_temp_frames(self):
        """Clean up temporary frame files."""
        for frame_path in self.frame_cache.values():
            try:
                if Path(frame_path).exists():
                    os.unlink(frame_path)
            except Exception:
                pass
        self.frame_cache.clear()
    
    def __del__(self):
        """Cleanup when object is destroyed."""
        self.cleanup_temp_frames()
