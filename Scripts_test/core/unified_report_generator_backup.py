import json, logging, sys, re, tempfile, os

from datetime import datetime
from pathlib import Path
from dataclasses import dataclass, field
from typing import Dict, List, Optional, Union
from concurrent.futures import ThreadPoolExecutor
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Cm, Inches, Pt
from pptx.enum.text import PP_ALIGN

try:
    import cv2
except ImportError:
    cv2 = None

logging.basicConfig(level=logging.INFO, format='%(message)s')
logger = logging.getLogger(__name__)

@dataclass
class MediaPair:
    """Universal media pair for all API types"""
    source_file: str
    source_path: Path
    api_type: str
    # Generated content
    generated_paths: List[Path] = field(default_factory=list)
    reference_paths: List[Path] = field(default_factory=list)
    # API-specific fields
    effect_name: str = ""
    category: str = ""
    prompt: str = ""
    # Video-specific (for Runway)
    source_video_path: Optional[Path] = None
    # Metadata
    metadata: Dict = field(default_factory=dict)
    ref_metadata: Dict = field(default_factory=dict)
    # Status
    failed: bool = False
    ref_failed: bool = False

    @property
    def primary_generated(self) -> Optional[Path]:
        """Get the primary generated file"""
        return self.generated_paths[0] if self.generated_paths else None

    @property
    def primary_reference(self) -> Optional[Path]:
        """Get the primary reference file"""
        return self.reference_paths[0] if self.reference_paths else None

class UnifiedReportGenerator:
    """
    Enhanced unified report generator with all missing functions from working versions
    """

    def __init__(self, api_name: str, config_file: str = None):
        self.api_name = api_name
        self.config_file = config_file or f"batch_{api_name}_config.json"
        self.config = {}
        self.report_definitions = {}

        # Caches for performance
        self._ar_cache = {}
        self._frame_cache = {}

        # Load configurations
        self._load_config()
        self._load_report_definitions()

    def _load_config(self):
        """Load API-specific configuration"""
        try:
            with open(self.config_file, 'r', encoding='utf-8') as f:
                self.config = json.load(f)
                logger.info(f"✓ Config loaded: {self.config_file}")
        except Exception as e:
            logger.error(f"❌ Config error: {e}")
            sys.exit(1)

    def _load_report_definitions(self):
        """Load report definitions from api_definitions.json"""
        try:
            with open("core/api_definitions.json", 'r', encoding='utf-8') as f:
                all_definitions = json.load(f)
                self.report_definitions = all_definitions.get(self.api_name, {}).get('report', {})
        except Exception as e:
            logger.warning(f"⚠️ Report definitions error: {e}, using defaults")
            self._set_default_report_definitions()

    def _set_default_report_definitions(self):
        """Set default report definitions"""
        self.report_definitions = {
            "enabled": True,
            "template_path": "templates/I2V-templates.pptx",
            "comparison_template_path": "templates/I2V-Comparison-Template.pptx",
            "output_directory": "/Users/ethanhsu/Desktop/GAI/Report",
            "use_comparison": self.api_name in ["kling", "nano_banana", "runway"]
        }

    # ================== MISSING FUNCTIONS FROM WORKING VERSIONS ==================

    def get_cmp_filename(self, folder1: str, folder2: str, model: str = '') -> str:
        """Generate comparison filename like '[0908] API Style1 vs Style2' (from nano_banana)"""
        def parse_name(name):
            m = re.match(r'(\d{4})\s+(.+)', Path(name).name if isinstance(name, str) else name)
            return m.groups() if m else (datetime.now().strftime("%m%d"), str(name))

        date, style1 = parse_name(folder1)
        _, style2 = parse_name(folder2)

        parts = [f"[{date}]"]
        if model and model.lower() not in style1.lower():
            parts.append(model)
        parts.append(f"{style1} vs {style2}")
        return ' '.join(parts)

    def calc_pos(self, path: Path, x: float, y: float, w: float, h: float) -> tuple:
        """Calculate positioned media with aspect ratio preservation (from nano_banana)"""
        try:
            with Image.open(path) as img:
                ar = img.width / img.height
            bw, bh = Cm(w), Cm(h)
            sw, sh = (bw, bw/ar) if ar > bw/bh else (bh*ar, bh)
            return Cm(x)+(bw-sw)/2, Cm(y)+(bh-sh)/2, sw, sh
        except:
            return Cm(x+0.5), Cm(y+0.5), Cm(w-1), Cm(h-1)

    def find_matching_video(self, base_name: str, video_files: dict) -> Optional[Path]:
        """Match image files to corresponding video files by stem (from kling optimized)"""
        if base_name in video_files:
            return video_files[base_name]

        for vname, vpath in video_files.items():
            if vname.startswith(base_name + '_'):
                return vpath
        return None

    def _add_hyperlink(self, para, prefix: str, link_text: str, url: str):
        """Simple hyperlink creation (from nano_banana)"""
        para.clear()
        r1 = para.add_run()
        r1.text = prefix
        r1.font.size = Inches(24/72)

        r2 = para.add_run()
        r2.text = link_text
        r2.font.size = Inches(24/72)
        r2.hyperlink.address = url

        para.alignment = PP_ALIGN.CENTER

    def get_aspect_ratio(self, path: Path, is_video: bool = False) -> float:
        """Enhanced aspect ratio calculation with caching (from working versions)"""
        filename = path.name.lower()

        # Quick detection from filename
        if '9_16' in filename or 'portrait' in filename:
            return 9/16
        if '1_1' in filename or 'square' in filename:
            return 1
        if '16_9' in filename or 'landscape' in filename:
            return 16/9

        # Check cache
        cache_key = str(path)
        if cache_key in self._ar_cache:
            return self._ar_cache[cache_key]

        try:
            if is_video and cv2:
                cap = cv2.VideoCapture(str(path))
                if cap.isOpened():
                    w, h = cap.get(3), cap.get(4)
                    cap.release()
                    if w > 0 and h > 0:
                        ar = w / h
                        self._ar_cache[cache_key] = ar
                        return ar
            else:
                with Image.open(path) as img:
                    ar = img.width / img.height
                    self._ar_cache[cache_key] = ar
                    return ar
        except Exception:
            pass

        # Default fallback
        return 16/9

    def extract_first_frame(self, video_path: Path) -> Optional[str]:
        """Extract first frame from video for poster (enhanced for all APIs)"""
        if not cv2:
            return None

        video_key = str(video_path)
        if video_key in self._frame_cache:
            return self._frame_cache[video_key]

        try:
            cap = cv2.VideoCapture(str(video_path))
            if not cap.isOpened():
                return None

            ret, frame = cap.read()
            cap.release()

            if not ret or frame is None:
                return None

            # Create temp file
            temp_dir = tempfile.gettempdir()
            frame_filename = f"frame_{self.api_name}_{video_path.stem}_{hash(video_key) % 10000}.jpg"
            frame_path = Path(temp_dir) / frame_filename

            cv2.imwrite(str(frame_path), frame)
            self._frame_cache[video_key] = str(frame_path)
            return str(frame_path)

        except Exception as e:
            logger.warning(f"Frame extraction failed: {e}")
            return None

    # ================== ENHANCED METADATA HANDLING ==================

    def _format_metadata_by_api(self, pair: MediaPair, use_comparison: bool = False) -> List[str]:
        """Format metadata based on API type with rich information"""
        metadata_lines = []

        if self.api_name == "kling":
            # Kling metadata: task_id, processing_time, model_version
            task_id = pair.metadata.get('task_id', 'N/A')
            proc_time = pair.metadata.get('processing_time_seconds', 'N/A')
            model = pair.metadata.get('model_version', 'v2.1')
            status = "SUCCESS" if pair.metadata.get('success', False) else "FAILED"

            metadata_lines = [
                f"Task: {task_id}",
                f"Model: Kling {model}",
                f"Time: {proc_time}s",
                f"Status: {status}"
            ]

            if use_comparison and pair.ref_metadata:
                ref_time = pair.ref_metadata.get('processing_time_seconds', 'N/A')
                metadata_lines.append(f"Ref Time: {ref_time}s")

        elif self.api_name == "nano_banana":
            # Nano Banana metadata: response_id, processing_time, images_generated
            resp_id = pair.metadata.get('response_id', 'N/A')
            proc_time = pair.metadata.get('processing_time_seconds', 'N/A')
            img_count = len(pair.generated_paths) if pair.generated_paths else 0

            metadata_lines = [
                f"File: {pair.source_file}",
                f"Response: {resp_id}",
                f"Images: {img_count} generated",
                f"Time: {proc_time}s"
            ]

        elif self.api_name == "runway":
            # Runway metadata: prompt, model, reference_image, processing_time
            prompt = pair.metadata.get('prompt', 'N/A')[:30] + "..." if len(pair.metadata.get('prompt', '')) > 30 else pair.metadata.get('prompt', 'N/A')
            model = pair.metadata.get('model', 'gen4_aleph')
            ref_img = pair.metadata.get('reference_image', 'N/A')
            proc_time = pair.metadata.get('processing_time_seconds', 'N/A')
            status = "✓" if pair.metadata.get('success', False) else "❌"

            metadata_lines = [
                f"Prompt: {prompt}",
                f"Ref: {ref_img}",
                f"Model: {model}",
                f"Time: {proc_time}s",
                f"Status: {status}"
            ]

        elif self.api_name == "vidu_effects":
            # Vidu Effects metadata: effect_name, category, processing_time
            proc_time = pair.metadata.get('processing_time_seconds', 'N/A')
            status = 'FAILED' if pair.failed else 'SUCCESS'

            metadata_lines = [
                f"{pair.source_file}",
                f"{pair.category} | {pair.effect_name}",
                f"{proc_time}s \\ {status}"
            ]

        elif self.api_name == "vidu_reference":
            # Vidu Reference metadata: effect_name, reference_count, aspect_ratio
            ref_count = pair.metadata.get('reference_count', 0)
            aspect_ratio = pair.metadata.get('detected_aspect_ratio', 'N/A')
            proc_time = pair.metadata.get('processing_time_seconds', 'N/A')
            status = 'FAILED' if pair.failed else 'SUCCESS'

            metadata_lines = [
                f"{pair.source_file}",
                f"{pair.effect_name} | +{ref_count} refs",
                f"{aspect_ratio} | {proc_time}s",
                f"{status}"
            ]

        else:
            # Default metadata format
            proc_time = pair.metadata.get('processing_time_seconds', 'N/A')
            status = "SUCCESS" if not pair.failed else "FAILED"
            metadata_lines = [
                f"Processing: {proc_time}s",
                f"Status: {status}"
            ]

            if pair.effect_name:
                metadata_lines.insert(0, f"Effect: {pair.effect_name}")

        return metadata_lines

    # ================== SIMPLIFIED SLIDE CREATION ==================

    def _create_title_slide_simple(self, ppt: Presentation, task: Dict, use_comparison: bool):
        """Simplified title slide creation using working version pattern"""
        if not ppt.slides:
            return

        # Get folder names for title generation
        folder_name = task.get('folder', Path(self.config.get('base_folder', '')).name)
        if isinstance(folder_name, str):
            folder_name = Path(folder_name).name

        api_display_names = {
            'kling': 'Kling 2.1',
            'nano_banana': 'Nano Banana',
            'runway': 'Runway',
            'vidu_effects': 'Vidu Effects',
            'vidu_reference': 'Vidu Reference'
        }
        api_display = api_display_names.get(self.api_name, self.api_name.title())

        # Simple title generation (from working versions)
        if use_comparison and task.get('reference_folder'):
            ref_name = Path(task['reference_folder']).name
            title = self.get_cmp_filename(folder_name, ref_name, api_display)
        else:
            m = re.match(r'(\d{4})\s+(.+)', folder_name)
            if m:
                date, style = m.groups()
                title = f"[{date}] {api_display}\n{style}"
            else:
                date = datetime.now().strftime("%m%d")
                title = f"[{date}] {api_display}\n{folder_name}"

        # SIMPLE APPROACH: Direct slides[0].shapes[0] assignment (from working versions)
        try:
            if ppt.slides and ppt.slides[0].shapes:
                ppt.slides[0].shapes[0].text_frame.text = title
                logger.info(f"✅ Title updated: {title}")
        except Exception:
            logger.warning("⚠️ Could not update title slide")

        # Add links using simple approach
        self._add_links_simple(ppt, task)

    def _add_links_simple(self, ppt: Presentation, task: Dict):
        """Simplified link addition (from working versions)"""
        if not ppt.slides:
            return

        slide = ppt.slides[0]

        # Find existing info box
        info_box = None
        for shape in slide.shapes:
            if (hasattr(shape, 'text_frame') and shape.text_frame.text and 
                any(x in shape.text_frame.text.lower() for x in ['design', 'testbed', 'source'])):
                info_box = shape
                break

        if not info_box:
            info_box = slide.shapes.add_textbox(Cm(5), Cm(13), Cm(20), Cm(4))

        info_box.text_frame.clear()

        # Get testbed URL from config
        testbed_url = self.config.get('testbed', f'http://192.168.4.3:8000/{self.api_name}/')

        links = [
            ("Design: ", "Link", task.get('design_link', '')),
            ("Testbed: ", testbed_url, testbed_url),
            ("Source + Video: ", "Link", task.get('source_video_link', ''))
        ]

        for i, (prefix, text, url) in enumerate(links):
            para = info_box.text_frame.paragraphs[0] if i == 0 else info_box.text_frame.add_paragraph()

            if url:
                self._add_hyperlink(para, prefix, text if "Testbed:" not in prefix else url, url)
            else:
                para.text = f"{prefix}{text}"
                para.font.size = Inches(24/72)
                para.alignment = PP_ALIGN.CENTER

    def _populate_slide_hybrid(self, slide, pair: MediaPair, index: int, use_comparison: bool):
        """Hybrid approach: try simple methods first, fallback to complex"""
        # Set title
        title = self._generate_slide_title(pair, index)

        # Find title placeholder
        title_placeholder = None
        for placeholder in slide.placeholders:
            if placeholder.placeholder_format.type == 1:
                title_placeholder = placeholder
                break

        if title_placeholder:
            title_placeholder.text = title
            if pair.failed or (use_comparison and pair.ref_failed):
                if title_placeholder.text_frame.paragraphs:
                    title_placeholder.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 0, 0)

        # Get content placeholders
        content_placeholders = sorted(
            [p for p in slide.placeholders if p.placeholder_format.type in {6, 7, 8, 13, 18, 19}],
            key=lambda x: getattr(x, 'left', 0)
        )

        if use_comparison and len(content_placeholders) >= 3:
            # 3-way comparison layout
            self._add_media_hybrid(slide, content_placeholders[0], pair.source_path)
            self._add_media_hybrid(slide, content_placeholders[1], pair.primary_generated, 
                                 is_video=self._is_video(pair.primary_generated),
                                 error=self._get_error_message(pair))
            self._add_media_hybrid(slide, content_placeholders[2], pair.primary_reference,
                                 is_video=self._is_video(pair.primary_reference), 
                                 error=self._get_ref_error_message(pair))
        elif len(content_placeholders) >= 2:
            # 2-way layout
            self._add_media_hybrid(slide, content_placeholders[0], pair.source_path)
            self._add_media_hybrid(slide, content_placeholders[1], pair.primary_generated,
                                 is_video=self._is_video(pair.primary_generated),
                                 error=self._get_error_message(pair))

        # Add enhanced metadata
        self._add_enhanced_metadata_box(slide, pair, use_comparison)

    def _add_media_hybrid(self, slide, placeholder, media_path: Optional[Path], 
                         is_video: bool = False, error: str = None):
        """Hybrid media addition: try insert_picture first, then positioning"""
        if not media_path or not media_path.exists():
            self._add_error_to_placeholder(slide, placeholder, error or "Media not found")
            return

        # TRY SIMPLE APPROACH FIRST (from working versions)
        if not is_video:
            try:
                placeholder.insert_picture(str(media_path))
                logger.info(f"✅ Direct insert: {media_path.name}")
                return
            except Exception:
                logger.info(f"⚠️ Direct insert failed, using positioning for: {media_path.name}")

        # FALLBACK TO POSITIONING (existing complex logic)
        p_left, p_top, p_width, p_height = placeholder.left, placeholder.top, placeholder.width, placeholder.height
        ar = self.get_aspect_ratio(media_path, is_video)

        if ar > p_width / p_height:
            scaled_w, scaled_h = p_width, p_width / ar
        else:
            scaled_h, scaled_w = p_height, p_height * ar

        final_left = p_left + (p_width - scaled_w) / 2
        final_top = p_top + (p_height - scaled_h) / 2

        # Remove placeholder
        placeholder._element.getparent().remove(placeholder._element)

        try:
            if is_video:
                poster_frame = self.extract_first_frame(media_path)
                if poster_frame and Path(poster_frame).exists():
                    slide.shapes.add_movie(str(media_path), final_left, final_top, scaled_w, scaled_h,
                                         poster_frame_image=poster_frame)
                else:
                    slide.shapes.add_movie(str(media_path), final_left, final_top, scaled_w, scaled_h)
            else:
                slide.shapes.add_picture(str(media_path), final_left, final_top, scaled_w, scaled_h)

            logger.info(f"✅ Positioned media: {media_path.name}")
        except Exception as e:
            logger.error(f"❌ Media add failed: {e}")
            self._add_error_box(slide, (final_left, final_top, scaled_w, scaled_h), f"Failed to load: {e}")

    def _add_enhanced_metadata_box(self, slide, pair: MediaPair, use_comparison: bool):
        """Add enhanced metadata box with API-specific formatting"""
        metadata_lines = self._format_metadata_by_api(pair, use_comparison)

        # Position based on API type and comparison mode
        if self.api_name in ["vidu_effects", "vidu_reference"]:
            # Compact metadata for vidu APIs
            meta_box = slide.shapes.add_textbox(Cm(2), Cm(16.5), Cm(8), Cm(2))
        else:
            # Standard metadata box
            meta_box = slide.shapes.add_textbox(Cm(2), Cm(16), Cm(12), Cm(3))

        meta_text = "\n".join(metadata_lines)
        meta_box.text_frame.text = meta_text
        meta_box.text_frame.word_wrap = True

        for para in meta_box.text_frame.paragraphs:
            para.font.size = Inches(10/72)

    # ================== EXISTING CORE FUNCTIONS (SIMPLIFIED) ==================

    def normalize_key(self, name: str) -> str:
        """Universal key normalization"""
        key = name.lower().replace(' ', '_')
        key = re.sub(r'[^a-z0-9_]', '', key)
        key = re.sub(r'_(effect|generated|output|result)$', '', key)
        return key.strip('_')

    def get_filename(self, folder: str, model: str = '', comparison_folder: str = '') -> str:
        """Universal filename generation with date parsing"""
        def parse_name(name):
            m = re.match(r'(?:\[(\d{4})\]|\[?(\d{4})\]?)\s*(.+)', name.strip())
            if m:
                return (m.group(1) or m.group(2), m.group(3).strip())
            return (datetime.now().strftime("%m%d"), name.strip())

        date, style = parse_name(folder)
        parts = [f"[{date}]"]

        if model and model.lower() not in style.lower():
            parts.append(model)

        if comparison_folder:
            _, ref_style = parse_name(comparison_folder)
            parts.append(f"{style} vs {ref_style}")
        else:
            parts.append(style)

        return ' '.join(parts)

    def process_batch(self, task: Dict) -> List[MediaPair]:
        """Universal batch processing for all API types"""
        if self.api_name in ["vidu_effects", "vidu_reference"]:
            return self._process_base_folder_structure(task)
        else:
            return self._process_task_folder_structure(task)

    def _process_task_folder_structure(self, task: Dict) -> List[MediaPair]:
        """Process APIs with individual task folders"""
        folder = Path(task['folder'])
        ref_folder = Path(task.get('reference_folder', '')) if task.get('reference_folder') else None
        use_comparison = task.get('use_comparison_template', False)

        # Define folder structure based on API
        if self.api_name == "nano_banana":
            folders = {
                'source': folder / 'Source',
                'generated': folder / 'Generated_Output',
                'metadata': folder / 'Metadata'
            }
        elif self.api_name == "runway":
            folders = {
                'source': folder / 'Source',
                'reference': folder / 'Reference',
                'generated': folder / 'Generated_Video',
                'metadata': folder / 'Metadata'
            }
        else:  # kling
            folders = {
                'source': folder / 'Source',
                'generated': folder / 'Generated_Video',
                'metadata': folder / 'Metadata'
            }

        # Reference folders for comparison
        ref_folders = {}
        if ref_folder and use_comparison:
            ref_folders = {
                'generated': ref_folder / ('Generated_Output' if self.api_name == "nano_banana" else 'Generated_Video'),
                'metadata': ref_folder / 'Metadata'
            }

        return self._create_media_pairs_from_folders(folders, ref_folders, task, use_comparison)

    def _process_base_folder_structure(self, task: Dict) -> List[MediaPair]:
        """Process APIs with base folder + effect subfolders"""
        base_folder = Path(self.config.get('base_folder', ''))
        pairs = []

        # Get effect names from config or discover
        if self.api_name == "vidu_effects":
            effects = [t.get('effect', '') for t in self.config.get('tasks', [])]
        else:  # vidu_reference
            effects = self._discover_effect_folders(base_folder)

        for effect in effects:
            if not effect:
                continue

            folders = {
                'source': base_folder / effect / 'Source',
                'generated': base_folder / effect / 'Generated_Video',
                'metadata': base_folder / effect / 'Metadata'
            }

            if self.api_name == "vidu_reference":
                folders['reference'] = base_folder / effect / 'Reference'

            effect_pairs = self._create_media_pairs_from_folders(folders, {}, {'effect': effect}, False)
            pairs.extend(effect_pairs)

        return pairs

    def _discover_effect_folders(self, base_folder: Path) -> List[str]:
        """Discover effect folders automatically"""
        if not base_folder.exists():
            return []

        return [f.name for f in base_folder.iterdir()
                if f.is_dir() and not f.name.startswith(('.', '_'))
                and (f / 'Source').exists()]

    def _create_media_pairs_from_folders(self, folders: Dict[str, Path], ref_folders: Dict[str, Path], 
                                       task: Dict, use_comparison: bool) -> List[MediaPair]:
        """Create MediaPair objects from folder contents with enhanced file matching"""
        pairs = []

        if not folders.get('source', Path()).exists():
            return pairs

        # Get source files
        source_exts = {'.jpg', '.jpeg', '.png', '.webp', '.bmp', '.tiff'}
        video_exts = {'.mp4', '.mov', '.avi', '.mkv'}

        if self.api_name == "runway":
            # Runway uses reference images as primary source
            source_files = {f.stem: f for f in folders.get('reference', Path()).iterdir()
                          if f.suffix.lower() in source_exts} if folders.get('reference', Path()).exists() else {}
        else:
            source_files = {f.stem: f for f in folders['source'].iterdir()
                          if f.suffix.lower() in source_exts}

        # Get generated files with enhanced matching
        generated_files = {}
        if folders.get('generated', Path()).exists():
            generated_dict = {}
            for f in folders['generated'].iterdir():
                if f.suffix.lower() in (video_exts | source_exts):
                    generated_dict[f.stem] = f

            # Use find_matching_video for better matching
            for source_stem in source_files.keys():
                matched = self.find_matching_video(source_stem, generated_dict)
                if matched:
                    generated_files[source_stem] = [matched]

        # Get metadata
        metadata_files = {}
        if folders.get('metadata', Path()).exists():
            for f in folders['metadata'].iterdir():
                if f.suffix.lower() == '.json':
                    key = f.stem.replace('_metadata', '')
                    key = self.normalize_key(key)
                    metadata_files[key] = f

        # Get reference files for comparison (similar logic for ref_folders)
        ref_generated_files = {}
        ref_metadata_files = {}
        if use_comparison and ref_folders:
            if ref_folders.get('generated', Path()).exists():
                ref_dict = {}
                for f in ref_folders['generated'].iterdir():
                    if f.suffix.lower() in (video_exts | source_exts):
                        ref_dict[f.stem] = f

                for source_stem in source_files.keys():
                    matched = self.find_matching_video(source_stem, ref_dict)
                    if matched:
                        ref_generated_files[source_stem] = [matched]

            if ref_folders.get('metadata', Path()).exists():
                for f in ref_folders['metadata'].iterdir():
                    if f.suffix.lower() == '.json':
                        key = f.stem.replace('_metadata', '')
                        key = self.normalize_key(key)
                        ref_metadata_files[key] = f

        # Create pairs
        for source_key, source_file in source_files.items():
            norm_key = self.normalize_key(source_key)

            # Load metadata
            metadata = self._load_metadata(metadata_files.get(norm_key))
            ref_metadata = self._load_metadata(ref_metadata_files.get(norm_key)) if use_comparison else {}

            # Get generated files
            generated_paths = generated_files.get(source_key, [])
            ref_generated_paths = ref_generated_files.get(source_key, []) if use_comparison else []

            # Handle special cases
            source_video_path = None
            if self.api_name == "runway" and folders.get('source', Path()).exists():
                # Find matching source video
                for f in folders['source'].iterdir():
                    if f.suffix.lower() in video_exts and self.normalize_key(f.stem) == norm_key:
                        source_video_path = f
                        break

            # Create MediaPair
            pair = MediaPair(
                source_file=source_file.name,
                source_path=source_file,
                api_type=self.api_name,
                generated_paths=generated_paths,
                reference_paths=ref_generated_paths,
                effect_name=task.get('effect', ''),
                category=task.get('category', ''),
                prompt=task.get('prompt', ''),
                source_video_path=source_video_path,
                metadata=metadata,
                ref_metadata=ref_metadata,
                failed=not generated_paths or not metadata.get('success', False),
                ref_failed=use_comparison and (not ref_generated_paths or not ref_metadata.get('success', False))
            )

            pairs.append(pair)

        return pairs

    def _load_metadata(self, metadata_file: Optional[Path]) -> Dict:
        """Load metadata from JSON file"""
        if not metadata_file or not metadata_file.exists():
            return {}

        try:
            with open(metadata_file, 'r', encoding='utf-8') as f:
                return json.load(f)
        except Exception:
            return {}

    def create_presentation(self, pairs: List[MediaPair], task: Dict) -> bool:
        """Universal presentation creation with enhanced features"""
        if not pairs:
            logger.warning("No media pairs to process")
            return False

        # Determine template and comparison mode
        use_comparison = self.report_definitions.get('use_comparison', False) and task.get('use_comparison_template', False)
        template_key = 'comparison_template_path' if use_comparison else 'template_path'
        template_path = (self.config.get(template_key) or
                        self.report_definitions.get(template_key,
                                                  'templates/I2V-Comparison-Template.pptx' if use_comparison else 'templates/I2V-templates.pptx'))

        # Load template
        try:
            ppt = Presentation(template_path) if Path(template_path).exists() else Presentation()
            template_loaded = Path(template_path).exists()
            logger.info(f"✓ Template loaded: {template_path}")
        except Exception as e:
            logger.warning(f"⚠️ Template load failed: {e}, using blank presentation")
            ppt = Presentation()
            template_loaded = False

        # Set slide dimensions
        ppt.slide_width, ppt.slide_height = Cm(33.87), Cm(19.05)

        # Create title slide using SIMPLE approach
        self._create_title_slide_simple(ppt, task, use_comparison)

        # Add content slides
        if self.api_name == "vidu_reference":
            self._create_sectioned_slides(ppt, pairs, template_loaded, use_comparison)
        else:
            self._create_standard_slides(ppt, pairs, template_loaded, use_comparison)

        # Save presentation
        return self._save_presentation(ppt, task, use_comparison)

    def _create_standard_slides(self, ppt: Presentation, pairs: List[MediaPair], template_loaded: bool, use_comparison: bool):
        """Create standard slides for most APIs"""
        for i, pair in enumerate(pairs, 1):
            self._create_content_slide(ppt, pair, i, template_loaded, use_comparison)

    def _create_sectioned_slides(self, ppt: Presentation, pairs: List[MediaPair], template_loaded: bool, use_comparison: bool):
        """Create sectioned slides for Vidu Reference"""
        from collections import defaultdict

        grouped = defaultdict(list)
        for pair in pairs:
            grouped[pair.effect_name].append(pair)

        for effect_name, effect_pairs in sorted(grouped.items()):
            # Create section slide
            try:
                section_slide = ppt.slides.add_slide(ppt.slide_layouts[1])
            except Exception:
                section_slide = ppt.slides.add_slide(ppt.slide_layouts[6])

            # Set section title
            for placeholder in section_slide.placeholders:
                if placeholder.placeholder_format.type == 1:
                    placeholder.text = effect_name
                    break

            # Create content slides
            for i, pair in enumerate(effect_pairs, 1):
                self._create_content_slide(ppt, pair, i, template_loaded, use_comparison)

    def _create_content_slide(self, ppt: Presentation, pair: MediaPair, index: int, template_loaded: bool, use_comparison: bool):
        """Create individual content slide using HYBRID approach"""
        # Try to use template layout
        try:
            if template_loaded and len(ppt.slides) >= 4:
                slide = ppt.slides.add_slide(ppt.slides[3].slide_layout)
                self._populate_slide_hybrid(slide, pair, index, use_comparison)
            else:
                slide = ppt.slides.add_slide(ppt.slide_layouts[6])
                self._populate_manual_slide(slide, pair, index, use_comparison)
        except Exception:
            slide = ppt.slides.add_slide(ppt.slide_layouts[6])
            self._populate_manual_slide(slide, pair, index, use_comparison)

    def _populate_manual_slide(self, slide, pair: MediaPair, index: int, use_comparison: bool):
        """Populate slide with manual positioning using calc_pos"""
        # Add title
        title = self._generate_slide_title(pair, index)
        title_box = slide.shapes.add_textbox(Cm(2), Cm(1), Cm(28), Cm(2))
        title_box.text_frame.text = title
        title_box.text_frame.paragraphs[0].font.size = Inches(18/72)

        if pair.failed or (use_comparison and pair.ref_failed):
            title_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 0, 0)

        # Position media using calc_pos from working versions
        positions = self._get_media_positions(use_comparison)

        # Add source media
        if pair.source_path.exists():
            x, y, w, h = self.calc_pos(pair.source_path, *positions[0])
            slide.shapes.add_picture(str(pair.source_path), x, y, w, h)

        # Add generated media
        if pair.primary_generated and pair.primary_generated.exists():
            x, y, w, h = self.calc_pos(pair.primary_generated, *positions[1])
            if self._is_video(pair.primary_generated):
                poster_frame = self.extract_first_frame(pair.primary_generated)
                if poster_frame and Path(poster_frame).exists():
                    slide.shapes.add_movie(str(pair.primary_generated), x, y, w, h, poster_frame_image=poster_frame)
                else:
                    slide.shapes.add_movie(str(pair.primary_generated), x, y, w, h)
            else:
                slide.shapes.add_picture(str(pair.primary_generated), x, y, w, h)
        else:
            self._add_error_box_simple(slide, positions[1], self._get_error_message(pair))

        # Add reference media for comparison
        if use_comparison and len(positions) > 2:
            if pair.primary_reference and pair.primary_reference.exists():
                x, y, w, h = self.calc_pos(pair.primary_reference, *positions[2])
                if self._is_video(pair.primary_reference):
                    poster_frame = self.extract_first_frame(pair.primary_reference)
                    if poster_frame and Path(poster_frame).exists():
                        slide.shapes.add_movie(str(pair.primary_reference), x, y, w, h, poster_frame_image=poster_frame)
                    else:
                        slide.shapes.add_movie(str(pair.primary_reference), x, y, w, h)
                else:
                    slide.shapes.add_picture(str(pair.primary_reference), x, y, w, h)
            else:
                self._add_error_box_simple(slide, positions[2], self._get_ref_error_message(pair))

        # Add enhanced metadata
        self._add_enhanced_metadata_box(slide, pair, use_comparison)

    # ================== HELPER FUNCTIONS ==================

    def _generate_slide_title(self, pair: MediaPair, index: int) -> str:
        """Generate slide title based on API type"""
        base_title = f"Generation {index}: {pair.source_file}"

        if self.api_name in ["vidu_effects", "vidu_reference"] and pair.effect_name:
            base_title = f"{pair.effect_name} #{index}: {pair.source_file}"

        if pair.failed:
            base_title += " ❌ FAILED"

        return base_title

    def _get_media_positions(self, use_comparison: bool) -> List[tuple]:
        """Get media positions based on comparison mode"""
        if use_comparison:
            return [(2.59, 3.26, 10, 10), (13, 3.26, 10, 10), (23.41, 3.26, 10, 10)]
        else:
            return [(2.59, 3.26, 12.5, 12.5), (18.78, 3.26, 12.5, 12.5)]

    def _add_error_to_placeholder(self, slide, placeholder, error_msg: str):
        """Add error message to placeholder"""
        p_left, p_top, p_width, p_height = placeholder.left, placeholder.top, placeholder.width, placeholder.height
        placeholder._element.getparent().remove(placeholder._element)
        self._add_error_box(slide, (p_left, p_top, p_width, p_height), error_msg)

    def _add_error_box(self, slide, position: tuple, message: str):
        """Add error box at position"""
        if isinstance(position[0], (int, float)):
            x, y, w, h = Cm(position[0]), Cm(position[1]), Cm(position[2]), Cm(position[3])
        else:
            x, y, w, h = position

        box = slide.shapes.add_textbox(x, y, w, h)
        box.text_frame.text = f"❌ GENERATION FAILED\n\n{message}"

        for para in box.text_frame.paragraphs:
            para.font.size = Inches(16/72)
            para.alignment = PP_ALIGN.CENTER
            para.font.color.rgb = RGBColor(255, 0, 0)

        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(255, 240, 240)
        box.line.color.rgb = RGBColor(255, 0, 0)
        box.line.width = Inches(0.02)

    def _add_error_box_simple(self, slide, position: tuple, message: str):
        """Simple error box for manual positioning"""
        x, y, w, h = position
        box = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
        box.text_frame.text = f"❌ FAILED\n\n{message}"

        for para in box.text_frame.paragraphs:
            para.font.size = Inches(14/72)
            para.alignment = PP_ALIGN.CENTER
            para.font.color.rgb = RGBColor(255, 0, 0)

        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(255, 240, 240)
        box.line.color.rgb = RGBColor(255, 0, 0)

    def _is_video(self, path: Optional[Path]) -> bool:
        """Check if file is a video"""
        if not path:
            return False
        return path.suffix.lower() in {'.mp4', '.mov', '.avi', '.mkv', '.webm'}

    def _get_error_message(self, pair: MediaPair) -> str:
        """Get appropriate error message for failed generation"""
        if pair.metadata and 'error' in pair.metadata:
            return pair.metadata['error']
        elif pair.metadata and 'error_message' in pair.metadata:
            return pair.metadata['error_message']
        return "Generation failed"

    def _get_ref_error_message(self, pair: MediaPair) -> str:
        """Get error message for failed reference"""
        if pair.ref_metadata and 'error' in pair.ref_metadata:
            return pair.ref_metadata['error']
        return "Reference not found"

    def _save_presentation(self, ppt: Presentation, task: Dict, use_comparison: bool) -> bool:
        """Save the presentation"""
        try:
            # Generate filename
            folder_name = task.get('folder', Path(self.config.get('base_folder', '')).name)
            if isinstance(folder_name, str):
                folder_name = Path(folder_name).name

            api_display_names = {
                'kling': 'Kling 2.1',
                'nano_banana': 'Nano Banana',
                'runway': 'Runway',
                'vidu_effects': 'Vidu Effects',
                'vidu_reference': 'Vidu Reference'
            }
            api_display = api_display_names.get(self.api_name, self.api_name.title())

            if use_comparison and task.get('reference_folder'):
                ref_name = Path(task['reference_folder']).name
                filename = self.get_cmp_filename(folder_name, ref_name, api_display)
            else:
                filename = self.get_filename(folder_name, api_display)

            # Get output directory
            output_dir = Path(self.config.get('output_directory',
                                            self.config.get('output', {}).get('directory',
                                                           self.report_definitions.get('output_directory', './'))))
            output_path = output_dir / f"{filename}.pptx"

            # Save
            ppt.save(str(output_path))
            logger.info(f"✅ Saved: {output_path}")
            return True

        except Exception as e:
            logger.error(f"❌ Save failed: {e}")
            return False

    def _cleanup_temp_frames(self):
        """Clean up temporary frame files"""
        for frame_path in self._frame_cache.values():
            try:
                if Path(frame_path).exists():
                    os.unlink(frame_path)
            except Exception:
                pass
        self._frame_cache.clear()

    def run(self) -> bool:
        """Main execution"""
        logger.info(f"🎬 Starting {self.api_name.title()} Report Generator")

        try:
            # Process tasks
            if self.api_name in ["vidu_effects", "vidu_reference"]:
                # Base folder structure - single task
                pairs = self.process_batch({})
                if pairs:
                    success = self.create_presentation(pairs, {})
                    logger.info(f"✓ Generated report with {len(pairs)} items")
                    return success
                else:
                    logger.warning("No media pairs found")
                    return False
            else:
                # Task folder structure - multiple tasks
                tasks = self.config.get('tasks', [])
                if not tasks:
                    logger.warning("No tasks found in configuration")
                    return False

                successful = 0
                for i, task in enumerate(tasks, 1):
                    logger.info(f"📁 Task {i}/{len(tasks)}: {Path(task['folder']).name}")
                    pairs = self.process_batch(task)

                    if pairs:
                        if self.create_presentation(pairs, task):
                            successful += 1
                        else:
                            logger.error(f"❌ Task {i} presentation failed")
                    else:
                        logger.warning(f"⚠️ Task {i} has no media pairs")

                logger.info(f"🎉 Generated {successful}/{len(tasks)} presentations")
                return successful > 0

        except Exception as e:
            logger.error(f"❌ Report generation failed: {e}")
            return False
        finally:
            self._cleanup_temp_frames()


# Factory function
def create_report_generator(api_name: str, config_file: str = None):
    """Factory function to create report generator"""
    return UnifiedReportGenerator(api_name, config_file)


# CLI entry point
def main():
    import argparse

    parser = argparse.ArgumentParser(description='Generate PowerPoint reports from API processing results')
    parser.add_argument('api_name', choices=['kling', 'nano_banana', 'vidu_effects', 'vidu_reference', 'runway'],
                       help='API type to generate report for')
    parser.add_argument('--config', '-c', help='Config file path (optional)')

    args = parser.parse_args()
    generator = create_report_generator(args.api_name, args.config)
    sys.exit(0 if generator.run() else 1)


if __name__ == "__main__":
    main()
